#!/usr/bin/env python3
"""Convert documents, ebooks, spreadsheets, slides, images, and text to LaTeX."""

from __future__ import annotations

import argparse
import base64
import csv
import html
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
import unicodedata
import urllib.error
import urllib.parse
import urllib.request
from collections import Counter
from collections.abc import Callable
from dataclasses import dataclass
from datetime import datetime, timezone
from difflib import SequenceMatcher
from email import policy
from email.parser import BytesParser
from pathlib import Path

import xlrd
from bs4 import BeautifulSoup
from charset_normalizer import from_bytes
from docx import Document
from ebooklib import ITEM_DOCUMENT, epub
from odf import draw as odf_draw
from odf import table as odf_table
from odf import text as odf_text
from odf.opendocument import load as load_odf
from odf.teletype import extractText as extract_odf_text
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from striprtf.striprtf import rtf_to_text

try:
    import fitz  # type: ignore
except Exception:  # noqa: BLE001
    fitz = None

try:
    from PIL import Image as _PILImage  # type: ignore
    from PIL import ImageOps as _PILImageOps  # type: ignore
except Exception:  # noqa: BLE001
    _PILImage = None
    _PILImageOps = None

try:
    from pillow_heif import register_heif_opener

    register_heif_opener()
except Exception:  # noqa: BLE001
    pass

try:
    import pytesseract  # type: ignore
except Exception:  # noqa: BLE001
    pytesseract = None

MATCH_MODE_EXACT = "exact"
MATCH_MODE_PERCENT = "percentage"

PAGE_FLOW_COMPACT = "compact"
PAGE_FLOW_SOURCE = "source_pages"
PAGE_SIZE_SOURCE = "source"
PAGE_SIZE_A4 = "a4"
PAGE_SIZE_LETTER = "letter"
PHOTO_KEEP = "keep"
PHOTO_DESCRIBE = "describe"

DEFAULT_OPENAI_COMPAT_ENDPOINT = "https://api.openai.com/v1/chat/completions"
DEFAULT_OLLAMA_ENDPOINT = "http://127.0.0.1:11434/api/chat"
APP_VERSION = "1.3.1"
GITHUB_LATEST_RELEASE_API = "https://api.github.com/repos/kkmmee94/book-to-latex/releases/latest"

DOCUMENT_LANGUAGES = {
    "auto": "Detect automatically",
    "eng": "English",
    "ara": "Arabic",
    "chi_sim": "Chinese (Simplified)",
    "chi_tra": "Chinese (Traditional)",
}

IMAGE_EXTENSIONS = {
    ".bmp",
    ".gif",
    ".heic",
    ".heif",
    ".jpeg",
    ".jpg",
    ".png",
    ".tif",
    ".tiff",
    ".webp",
}
DOCUMENT_EXTENSIONS = {
    ".docx",
    ".epub",
    ".htm",
    ".html",
    ".odt",
    ".ods",
    ".odp",
    ".rtf",
    ".eml",
}
SPREADSHEET_EXTENSIONS = {".csv", ".tsv", ".xls", ".xlsx", ".xlsm"}
PRESENTATION_EXTENSIONS = {".pptx"}
LEGACY_OFFICE_EXTENSIONS = {".doc", ".dot", ".ppt", ".pps", ".xls", ".wps"}
KNOWN_TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".rst",
    ".tex",
    ".json",
    ".jsonl",
    ".xml",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".conf",
    ".log",
    ".sql",
    ".py",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".css",
    ".scss",
    ".java",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".cs",
    ".go",
    ".rs",
    ".sh",
    ".ps1",
    ".bat",
}
SUPPORTED_INPUT_EXTENSIONS = (
    {".pdf"}
    | IMAGE_EXTENSIONS
    | DOCUMENT_EXTENSIONS
    | SPREADSHEET_EXTENSIONS
    | PRESENTATION_EXTENSIONS
    | LEGACY_OFFICE_EXTENSIONS
    | KNOWN_TEXT_EXTENSIONS
)

DEFAULT_LATEX_STYLE_GUIDE = """
- Preserve every word, number, formula, symbol, source typo, and ordering unless correction is explicitly requested.
- Return LaTeX body content only; never emit a document class, package list, or Markdown fence.
- Use amsmath/mathtools environments such as align*, cases, and gathered for mathematics.
- Use semantic LaTeX rather than visual spacing hacks. Use booktabs conventions for tables and listings for code/console output.
- Keep enumeration labels structural with enumerate/itemize rather than manually typed labels.
- Never guess an unreadable formula or numeric value; preserve the extracted source and let the review report flag uncertainty.
""".strip()

ProgressCallback = Callable[[int, int, str], None]
LogCallback = Callable[[str], None]
CancelCallback = Callable[[], bool]


class ConversionCancelled(RuntimeError):
    """Raised when a caller asks an in-progress conversion to stop."""


def _resource_path(*parts: str) -> Path:
    """Locate a bundled resource in source checkouts and PyInstaller apps."""
    bundle_root = Path(getattr(sys, "_MEIPASS", Path(__file__).resolve().parent))
    return bundle_root.joinpath(*parts)


def _bundled_tessdata_dir() -> Path | None:
    candidates = [
        _resource_path("assets", "tessdata"),
        Path(sys.prefix) / "share" / "book-to-latex" / "tessdata",
    ]
    return next((directory for directory in candidates if directory.is_dir()), None)


def _user_tessdata_dir() -> Path:
    if os.name == "nt":
        root = Path(os.environ.get("LOCALAPPDATA", Path.home() / "AppData/Local"))
        return root / "BookToLatex" / "tessdata"
    if sys.platform == "darwin":
        return Path.home() / "Library" / "Application Support" / "BookToLatex" / "tessdata"
    root = Path(os.environ.get("XDG_DATA_HOME", Path.home() / ".local/share"))
    return root / "book-to-latex" / "tessdata"


def ensure_ocr_language(
    language: str,
    log_callback: LogCallback | None = None,
) -> list[str]:
    """Download missing official Tesseract language data on demand."""
    codes = [code for code in (language or "eng").split("+") if code]
    allowed = set(DOCUMENT_LANGUAGES) - {"auto"}
    invalid = [code for code in codes if code not in allowed]
    if invalid:
        raise ValueError(f"Unsupported OCR language code: {', '.join(invalid)}")

    installed: set[str] = set()
    tesseract_path = _find_tesseract()
    if pytesseract is not None and tesseract_path:
        try:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
            installed = set(pytesseract.get_languages(config=""))
        except Exception:  # noqa: BLE001
            installed = set()
    bundled = _bundled_tessdata_dir()
    user_dir = _user_tessdata_dir()
    ready: list[str] = []
    for code in codes:
        if code in installed or (bundled and (bundled / f"{code}.traineddata").is_file()):
            ready.append(code)
            continue
        destination = user_dir / f"{code}.traineddata"
        if destination.is_file() and destination.stat().st_size > 100_000:
            ready.append(code)
            continue
        user_dir.mkdir(parents=True, exist_ok=True)
        url = f"https://raw.githubusercontent.com/tesseract-ocr/tessdata_fast/main/{code}.traineddata"
        temporary = destination.with_suffix(".traineddata.download")
        _log(log_callback, f"Installing OCR language data: {DOCUMENT_LANGUAGES.get(code, code)}")
        try:
            request = urllib.request.Request(
                url,
                headers={"User-Agent": f"book-to-latex/{APP_VERSION}"},
            )
            with urllib.request.urlopen(request, timeout=120) as response:
                data = response.read(100 * 1024 * 1024 + 1)
            if not 100_000 < len(data) <= 100 * 1024 * 1024:
                raise RuntimeError("Downloaded OCR language file has an unexpected size")
            temporary.write_bytes(data)
            temporary.replace(destination)
        except Exception as exc:  # noqa: BLE001
            if temporary.is_file():
                temporary.unlink()
            raise RuntimeError(
                f"Could not install {DOCUMENT_LANGUAGES.get(code, code)} OCR data: {exc}"
            ) from exc
        ready.append(code)
    return ready


def _tesseract_config(language: str) -> str:
    """Use bundled OCR data when the selected language is supplied by the app."""
    language_code = (language or "eng").split("+")[0]
    directories = [_user_tessdata_dir(), _bundled_tessdata_dir()]
    directory = next(
        (
            candidate
            for candidate in directories
            if candidate and (candidate / f"{language_code}.traineddata").is_file()
        ),
        None,
    )
    if directory:
        directory_text = str(directory)
        if os.name == "nt":
            # pytesseract uses non-POSIX shlex on Windows, which preserves quote
            # characters. Prefer an unquoted 8.3 path when whitespace is present.
            if any(char.isspace() for char in directory_text):
                try:
                    import ctypes

                    buffer = ctypes.create_unicode_buffer(32768)
                    length = ctypes.windll.kernel32.GetShortPathNameW(  # type: ignore[attr-defined]
                        directory_text,
                        buffer,
                        len(buffer),
                    )
                    if length:
                        directory_text = buffer.value
                except Exception:  # noqa: BLE001
                    pass
            return f"--tessdata-dir {directory_text}"
        return f'--tessdata-dir "{directory_text}"'
    return ""


@dataclass
class PageOutput:
    page_no: int
    source_text: str
    llm_text: str
    latex_body: str
    match_ratio: float
    uncertain: bool
    image_path: Path | None = None
    warnings: str | None = None
    used_fallback: bool = False
    number_match: bool = True
    missing_numbers: str = ""
    extra_numbers: str = ""


UNICODE_TO_LATEX = {
    "\u2018": "'",
    "\u2019": "'",
    "\u201c": '"',
    "\u201d": '"',
    "\u2013": "--",
    "\u2014": "---",
    "\u2026": "...",
    "\u00b0": "\\textdegree{}",
    "\u00d7": "$\\times$",
    "\u00f7": "$\\div$",
    "\u2211": "$\\sum$",
    "\u2202": "$\\partial$",
    "\u220f": "$\\prod$",
    "\u221e": "$\\infty$",
    "\u2265": "$\\geq$",
    "\u2264": "$\\leq$",
    "\u2260": "$\\neq$",
    "\u2248": "$\\approx$",
    "\u221d": "$\\propto$",
    "\u221a": "$\\surd$",
    "\u2022": "\\textbullet{}",
    "\u00b1": "$\\pm$",
}

LATEX_SPECIAL_CHARS = {
    "\\": r"\textbackslash{}",
    "&": r"\&",
    "%": r"\%",
    "#": r"\#",
    "_": r"\_",
    "{": r"\{",
    "}": r"\}",
    "^": r"\textasciicircum{}",
    "~": r"\textasciitilde{}",
    "$": r"\$",
}


def _log(log_callback: LogCallback | None, message: str) -> None:
    if log_callback:
        log_callback(message)


def _progress(progress_callback: ProgressCallback | None, current: int, total: int, message: str) -> None:
    if progress_callback:
        progress_callback(current, total, message)


def _sanitize_unicode(text: str, strict_mode: bool) -> str:
    if not strict_mode:
        return text
    out = text
    for source, replacement in UNICODE_TO_LATEX.items():
        out = out.replace(source, replacement)
    return out


def _escape_latex(text: str, strict_mode: bool = True) -> str:
    """Escape plain text once without re-escaping replacement commands."""
    rendered: list[str] = []
    for character in text:
        if strict_mode and character in UNICODE_TO_LATEX:
            rendered.append(UNICODE_TO_LATEX[character])
        else:
            rendered.append(LATEX_SPECIAL_CHARS.get(character, character))
    return "".join(rendered)


def _strip_markdown_fences(text: str) -> str:
    content = text.strip()
    fenced = re.fullmatch(r"```(?:latex|tex)?\s*\n?(.*?)\n?```", content, flags=re.DOTALL | re.IGNORECASE)
    if fenced:
        return fenced.group(1).strip()
    return content


def _extract_body_only(text: str) -> str:
    if "\\begin{document}" in text and "\\end{document}" in text:
        body = text.split("\\begin{document}", 1)[1]
        return body.split("\\end{document}", 1)[0].strip()
    return text.strip()


def _paragraphs_from_text(text: str, preserve_layout: bool, strict_mode: bool = True) -> str:
    normalized = text.replace("\r\n", "\n").strip()
    if not normalized:
        return ""

    if preserve_layout:
        lines = normalized.split("\n")
        rendered: list[str] = []
        for line in lines:
            if not line.strip():
                rendered.append("\\par")
            else:
                rendered.append(f"{_escape_latex(line, strict_mode=strict_mode)}\\\\")
        return "\n".join(rendered)

    paras = re.split(r"\n{2,}", normalized)
    rendered = [
        _escape_latex(paragraph, strict_mode=strict_mode)
        for paragraph in paras
        if paragraph.strip()
    ]
    return "\n\n\\par\n".join(rendered)


def _fallback_latex(text: str, preserve_layout: bool = True, strict_mode: bool = True) -> str:
    return _paragraphs_from_text(
        text,
        preserve_layout=preserve_layout,
        strict_mode=strict_mode,
    )


def _clean_latex_output(text: str, strict_mode: bool = True) -> str:
    text = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL | re.IGNORECASE)
    clean = _strip_markdown_fences(text)
    clean = _extract_body_only(clean)
    clean = _sanitize_unicode(clean, strict_mode=strict_mode)
    return clean.strip()


def _post_validate_latex(text: str) -> tuple[str, list[str]]:
    issues: list[str] = []
    cleaned = text

    if "```" in cleaned:
        cleaned = cleaned.replace("```", "")
        issues.append("Removed markdown fences")

    opens = cleaned.count("{")
    closes = cleaned.count("}")
    if opens != closes:
        issues.append(f"Unbalanced braces detected ({opens} open, {closes} close)")

    if "\\end{document}" in cleaned and "\\begin{document}" not in cleaned:
        issues.append("Unexpected end{document} without begin{document}")

    if "\\documentclass" in cleaned or "\\usepackage" in cleaned:
        issues.append("Unexpected document preamble")

    alignment_environments = (
        "align",
        "align*",
        "aligned",
        "array",
        "cases",
        "matrix",
        "pmatrix",
        "bmatrix",
        "vmatrix",
        "Vmatrix",
        "split",
        "tabular",
        "tabular*",
    )
    inside_alignment = False
    repaired_lines: list[str] = []
    repaired_ampersand = False
    repaired_percent = False
    repaired_hash = False
    for line in cleaned.splitlines():
        if any(f"\\begin{{{environment}}}" in line for environment in alignment_environments):
            inside_alignment = True
        if not inside_alignment:
            line, ampersands = re.subn(r"(?<!\\)&", r"\\&", line)
            repaired_ampersand = repaired_ampersand or ampersands > 0
        line, percents = re.subn(r"(?<!\\)%", r"\\%", line)
        line, hashes = re.subn(r"(?<!\\)#", r"\\#", line)
        repaired_percent = repaired_percent or percents > 0
        repaired_hash = repaired_hash or hashes > 0
        repaired_lines.append(line)
        if any(f"\\end{{{environment}}}" in line for environment in alignment_environments):
            inside_alignment = False
    cleaned = "\n".join(repaired_lines)
    if repaired_ampersand:
        issues.append("Escaped an unprotected ampersand in ordinary text")
    if repaired_percent:
        issues.append("Escaped an unprotected percent sign")
    if repaired_hash:
        issues.append("Escaped an unprotected hash sign")

    cleaned = cleaned.strip()
    return cleaned, issues


def _normalize_plain(text: str) -> str:
    """Normalize insignificant layout while preserving case and punctuation."""
    plain = unicodedata.normalize("NFKC", text)
    typography = str.maketrans(
        {
            "\u2018": "'",
            "\u2019": "'",
            "\u201c": '"',
            "\u201d": '"',
            "\u2013": "--",
            "\u2014": "---",
            "\u2026": "...",
        }
    )
    plain = plain.translate(typography)
    plain = re.sub(r"\s+", " ", plain)
    return plain.strip()


def _latex_to_plain(text: str) -> str:
    """Extract visible text from generated LaTeX for the review comparison."""
    plain = re.sub(r"(?<!\\)%.*$", "", text, flags=re.MULTILINE)
    plain = re.sub(
        r"\\includegraphics(?:\[[^\]]*\])?\s*\{[^{}]*\}",
        " ",
        plain,
        flags=re.DOTALL,
    )
    plain = re.sub(r"\\(?:begin|end)\s*\{[^{}]*\}", " ", plain)

    replacements = {
        r"\textbackslash{}": "\ue003",
        r"\textasciicircum{}": "^",
        r"\textasciitilde{}": "~",
        r"\textdegree{}": "\u00b0",
        r"\textbullet{}": "\u2022",
        r"\times": "\u00d7",
        r"\div": "\u00f7",
        r"\sum": "\u2211",
        r"\partial": "\u2202",
        r"\prod": "\u220f",
        r"\infty": "\u221e",
        r"\geq": "\u2265",
        r"\leq": "\u2264",
        r"\neq": "\u2260",
        r"\approx": "\u2248",
        r"\propto": "\u221d",
        r"\surd": "\u221a",
        r"\pm": "\u00b1",
        r"\&": "&",
        r"\%": "%",
        r"\#": "#",
        r"\_": "_",
        r"\{": "\ue000",
        r"\}": "\ue001",
        r"\$": "\ue002",
    }
    for latex, visible in sorted(replacements.items(), key=lambda item: len(item[0]), reverse=True):
        plain = plain.replace(latex, visible)

    formatting = re.compile(
        r"\\(?:textbf|textit|texttt|textrm|textsf|textsc|emph|underline|mbox)\*?\s*\{([^{}]*)\}"
    )
    while formatting.search(plain):
        plain = formatting.sub(r"\1", plain)

    plain = plain.replace(r"\\", " ")
    plain = plain.replace(r"\par", " ")
    plain = re.sub(r"\\[a-zA-Z@]+\*?(?:\[[^\]]*\])?", " ", plain)
    plain = re.sub(r"[{}$]", " ", plain)
    plain = (
        plain.replace("\ue000", "{")
        .replace("\ue001", "}")
        .replace("\ue002", "$")
        .replace("\ue003", "\\")
    )
    return _normalize_plain(plain)


def _match_similarity(a: str, b: str) -> float:
    if not a and not b:
        return 1.0
    return SequenceMatcher(None, a, b).ratio()


def _numeric_tokens(text: str) -> list[str]:
    return re.findall(
        r"(?<![\w.])[-+]?(?:\d{1,3}(?:,\d{3})+|\d+)(?:\.\d+)?(?:[eE][-+]?\d+)?%?",
        text,
    )


def _counter_difference(left: Counter[str], right: Counter[str], limit: int = 40) -> list[str]:
    output: list[str] = []
    for value, count in (left - right).most_common(limit):
        output.extend([value] * count)
    return output


def _numeric_comparison(source: str, generated: str) -> dict[str, object]:
    source_numbers = Counter(_numeric_tokens(source))
    generated_numbers = Counter(_numeric_tokens(generated))
    source_digits = "".join(
        character for token in _numeric_tokens(source) for character in token if character.isdigit()
    )
    generated_digits = "".join(
        character for token in _numeric_tokens(generated) for character in token if character.isdigit()
    )
    token_match = source_numbers == generated_numbers
    digit_match = source_digits == generated_digits
    return {
        "source_number_count": sum(source_numbers.values()),
        "generated_number_count": sum(generated_numbers.values()),
        "number_tokens_match_exactly": token_match,
        "numeric_digits_match": digit_match,
        "numbers_match": token_match or digit_match,
        "missing_numbers": _counter_difference(source_numbers, generated_numbers),
        "extra_numbers": _counter_difference(generated_numbers, source_numbers),
    }


def _verification_from_text(source: str, generated: str) -> dict[str, object]:
    source = _normalize_plain(source)
    generated = _normalize_plain(generated)
    source_words = re.findall(r"[^\W_]+(?:['’-][^\W_]+)*", source.casefold(), flags=re.UNICODE)
    generated_words = re.findall(
        r"[^\W_]+(?:['’-][^\W_]+)*", generated.casefold(), flags=re.UNICODE
    )
    source_word_counter = Counter(source_words)
    generated_word_counter = Counter(generated_words)
    verification = {
        "source_word_count": len(source_words),
        "generated_word_count": len(generated_words),
        "word_count_ratio": (
            len(generated_words) / len(source_words) if source_words else 1.0
        ),
        "text_similarity": _match_similarity(_normalize_plain(source), _normalize_plain(generated)),
        "missing_word_tokens": _counter_difference(source_word_counter, generated_word_counter),
        "extra_word_tokens": _counter_difference(generated_word_counter, source_word_counter),
    }
    verification.update(_numeric_comparison(source, generated))
    return verification


def _verification_summary(page_outputs: list[PageOutput]) -> dict[str, object]:
    source = "\n".join(page.source_text for page in page_outputs)
    generated = "\n".join(_latex_to_plain(page.latex_body) for page in page_outputs)
    return _verification_from_text(source, generated)


def _verify_compiled_pdf(
    pdf_path: Path,
    page_outputs: list[PageOutput],
) -> dict[str, object]:
    pages, extraction_engine = _read_pdf_pages(pdf_path)
    visible_lines = []
    for line in "\n".join(pages).splitlines():
        stripped = line.strip()
        if not stripped or stripped.isdigit() or stripped == "Converted manuscript":
            continue
        visible_lines.append(stripped)
    source = "\n".join(page.source_text for page in page_outputs)
    result = _verification_from_text(source, "\n".join(visible_lines))
    result["extraction_engine"] = extraction_engine
    result["compiled_pdf_pages"] = len(pages)
    return result


def _text_to_pages(text: str, lines_per_page: int) -> list[str]:
    lines_per_page = max(lines_per_page, 1)
    lines = text.replace("\r\n", "\n").split("\n")
    pages: list[str] = []
    for i in range(0, len(lines), lines_per_page):
        block = "\n".join(lines[i : i + lines_per_page]).strip()
        if block or not pages:
            pages.append(block)
    if not pages:
        pages.append("")
    return pages


def _decode_text_bytes(data: bytes, file_name: str) -> str:
    if not data:
        return ""
    if data.startswith((b"\xff\xfe\x00\x00", b"\x00\x00\xfe\xff")):
        return data.decode("utf-32")
    if data.startswith((b"\xff\xfe", b"\xfe\xff")):
        return data.decode("utf-16")
    if b"\x00" in data:
        raise ValueError(
            f"{file_name} appears to be an unsupported binary file rather than readable text"
        )
    try:
        text = data.decode("utf-8-sig")
    except UnicodeDecodeError:
        match = from_bytes(data).best()
        if match is None:
            raise ValueError(f"Could not detect the text encoding in {file_name}") from None
        text = str(match)
    if text:
        readable = sum(character.isprintable() or character in "\r\n\t" for character in text)
        if readable / len(text) < 0.85 or "\x00" in text:
            raise ValueError(
                f"{file_name} appears to be an unsupported binary file rather than readable text"
            )
    return text


def _read_text_pages(file_path: Path, lines_per_page: int) -> list[str]:
    return _text_to_pages(
        _decode_text_bytes(file_path.read_bytes(), file_path.name),
        lines_per_page,
    )


def _read_docx_pages(file_path: Path, lines_per_page: int) -> list[str]:
    document = Document(str(file_path))
    lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table_number, table in enumerate(document.tables, start=1):
        lines.append(f"[Table {table_number}]")
        for row in table.rows:
            values = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            lines.append("\t".join(values))
    return _text_to_pages("\n".join(lines), lines_per_page)


def _read_pptx_pages(file_path: Path) -> list[str]:
    presentation = Presentation(str(file_path))
    pages: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        lines = [f"[Slide {slide_number}]"]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                lines.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    lines.append("\t".join(cell.text.strip() for cell in row.cells))
        pages.append("\n".join(lines).strip())
    return pages or [""]


def _sheet_rows_to_pages(
    sheet_name: str,
    rows: list[list[object]],
    lines_per_page: int,
) -> list[str]:
    rendered_rows = [
        "\t".join("" if value is None else str(value) for value in row).rstrip()
        for row in rows
    ]
    rendered_rows = [row for row in rendered_rows if row.strip()]
    chunks = _text_to_pages("\n".join(rendered_rows), max(lines_per_page - 1, 1))
    return [f"[Sheet: {sheet_name}]\n{chunk}".strip() for chunk in chunks]


def _read_xlsx_pages(file_path: Path, lines_per_page: int) -> list[str]:
    workbook = load_workbook(filename=file_path, read_only=True, data_only=True)
    pages: list[str] = []
    try:
        for sheet in workbook.worksheets:
            rows = [list(row) for row in sheet.iter_rows(values_only=True)]
            pages.extend(_sheet_rows_to_pages(sheet.title, rows, lines_per_page))
    finally:
        workbook.close()
    return pages or [""]


def _read_xls_pages(file_path: Path, lines_per_page: int) -> list[str]:
    workbook = xlrd.open_workbook(str(file_path), on_demand=True)
    pages: list[str] = []
    try:
        for sheet in workbook.sheets():
            rows = [sheet.row_values(row_number) for row_number in range(sheet.nrows)]
            pages.extend(_sheet_rows_to_pages(sheet.name, rows, lines_per_page))
    finally:
        workbook.release_resources()
    return pages or [""]


def _read_html_pages(file_path: Path, lines_per_page: int) -> list[str]:
    markup = _decode_text_bytes(file_path.read_bytes(), file_path.name)
    soup = BeautifulSoup(markup, "html.parser")
    for removable in soup(["script", "style", "noscript"]):
        removable.decompose()
    return _text_to_pages(soup.get_text("\n", strip=True), lines_per_page)


def _read_epub_pages(file_path: Path, lines_per_page: int) -> list[str]:
    book = epub.read_epub(str(file_path))
    pages: list[str] = []
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        title = soup.title.get_text(" ", strip=True) if soup.title else item.get_name()
        content = soup.get_text("\n", strip=True)
        for chunk in _text_to_pages(content, max(lines_per_page - 1, 1)):
            pages.append(f"[Section: {title}]\n{chunk}".strip())
    return pages or [""]


def _read_odf_pages(file_path: Path, lines_per_page: int) -> list[str]:
    document = load_odf(str(file_path))
    extension = file_path.suffix.lower()
    if extension == ".ods":
        pages: list[str] = []
        for sheet in document.getElementsByType(odf_table.Table):
            sheet_name = sheet.getAttribute("name") or "Sheet"
            rows: list[list[object]] = []
            for row in sheet.getElementsByType(odf_table.TableRow):
                values = [
                    extract_odf_text(cell).replace("\n", " ").strip()
                    for cell in row.getElementsByType(odf_table.TableCell)
                ]
                rows.append(values)
            pages.extend(_sheet_rows_to_pages(sheet_name, rows, lines_per_page))
        return pages or [""]
    if extension == ".odp":
        pages = []
        for slide_number, slide in enumerate(document.getElementsByType(odf_draw.Page), start=1):
            text = "\n".join(
                extract_odf_text(node).strip()
                for node in slide.getElementsByType(odf_text.P)
                if extract_odf_text(node).strip()
            )
            pages.append(f"[Slide {slide_number}]\n{text}".strip())
        return pages or [""]
    paragraphs = [
        extract_odf_text(node).strip()
        for node_type in (odf_text.H, odf_text.P)
        for node in document.getElementsByType(node_type)
        if extract_odf_text(node).strip()
    ]
    return _text_to_pages("\n".join(paragraphs), lines_per_page)


def _read_email_pages(file_path: Path, lines_per_page: int) -> list[str]:
    message = BytesParser(policy=policy.default).parsebytes(file_path.read_bytes())
    lines = [
        f"Subject: {message.get('subject', '')}",
        f"From: {message.get('from', '')}",
        f"To: {message.get('to', '')}",
        f"Date: {message.get('date', '')}",
        "",
    ]
    bodies: list[str] = []
    parts = message.walk() if message.is_multipart() else [message]
    for part in parts:
        if part.get_content_disposition() == "attachment":
            continue
        if part.get_content_type() == "text/plain":
            content = part.get_content()
            if isinstance(content, str):
                bodies.append(content)
        elif part.get_content_type() == "text/html" and not bodies:
            content = part.get_content()
            if isinstance(content, str):
                bodies.append(BeautifulSoup(content, "html.parser").get_text("\n", strip=True))
    return _text_to_pages("\n".join(lines + bodies), lines_per_page)


def _find_libreoffice() -> str | None:
    found = shutil.which("soffice")
    if found:
        return found
    if os.name == "nt":
        for program_root in (os.environ.get("ProgramFiles"), os.environ.get("ProgramFiles(x86)")):
            if program_root:
                candidate = Path(program_root) / "LibreOffice/program/soffice.exe"
                if candidate.is_file():
                    return str(candidate)
    return None


def _read_legacy_office_pages(file_path: Path) -> list[str]:
    libreoffice = _find_libreoffice()
    if not libreoffice:
        raise RuntimeError(
            f"{file_path.suffix.upper()} is a legacy Office format. Install LibreOffice or save it as a modern DOCX, PPTX, or XLSX file."
        )
    with tempfile.TemporaryDirectory() as temporary:
        completed = subprocess.run(
            [libreoffice, "--headless", "--convert-to", "pdf", "--outdir", temporary, str(file_path)],
            capture_output=True,
            text=True,
            timeout=180,
            check=False,
        )
        converted = Path(temporary) / f"{file_path.stem}.pdf"
        if completed.returncode != 0 or not converted.is_file():
            detail = (completed.stderr or completed.stdout).strip()
            raise RuntimeError(f"LibreOffice could not read {file_path.name}: {detail}")
        return _read_pdf_with_pypdf(converted)


def _read_document_pages(file_path: Path, lines_per_page: int) -> tuple[list[str], str]:
    extension = file_path.suffix.lower()
    if extension == ".docx":
        return _read_docx_pages(file_path, lines_per_page), "Word document"
    if extension == ".pptx":
        return _read_pptx_pages(file_path), "PowerPoint presentation"
    if extension in {".xlsx", ".xlsm"}:
        return _read_xlsx_pages(file_path, lines_per_page), "Excel workbook"
    if extension == ".xls":
        return _read_xls_pages(file_path, lines_per_page), "Excel workbook"
    if extension in {".csv", ".tsv"}:
        return _read_text_pages(file_path, lines_per_page), "delimited data"
    if extension in {".html", ".htm"}:
        return _read_html_pages(file_path, lines_per_page), "web page"
    if extension == ".epub":
        return _read_epub_pages(file_path, lines_per_page), "EPUB ebook"
    if extension in {".odt", ".ods", ".odp"}:
        return _read_odf_pages(file_path, lines_per_page), "OpenDocument file"
    if extension == ".rtf":
        source = _decode_text_bytes(file_path.read_bytes(), file_path.name)
        return _text_to_pages(rtf_to_text(source), lines_per_page), "Rich Text document"
    if extension == ".eml":
        return _read_email_pages(file_path, lines_per_page), "email message"
    if extension in LEGACY_OFFICE_EXTENSIONS:
        return _read_legacy_office_pages(file_path), "legacy Office document"
    return _read_text_pages(file_path, lines_per_page), "text file"


def _read_pdf_with_pypdf(file_path: Path) -> list[str]:
    reader = PdfReader(str(file_path))
    pages: list[str] = []
    for page in reader.pages:
        pages.append((page.extract_text() or "").strip())
    return pages


def _ensure_fitz() -> None:
    if fitz is None:
        raise RuntimeError(
            "PyMuPDF (fitz) is required for PDF handling. Install with `pip install PyMuPDF`."
        )


def _ensure_pillow() -> None:
    if _PILImage is None:
        raise RuntimeError("Pillow is required for OCR. Install with `pip install Pillow`.")


def _find_tesseract() -> str | None:
    found = shutil.which("tesseract")
    if found:
        return found
    if os.name == "nt":
        candidates = [
            Path(os.environ.get("ProgramFiles", "C:/Program Files")) / "Tesseract-OCR/tesseract.exe",
            Path(os.environ.get("LOCALAPPDATA", "")) / "Programs/Tesseract-OCR/tesseract.exe",
        ]
        for candidate in candidates:
            if candidate.is_file():
                return str(candidate)
    return None


def _find_ollama() -> str | None:
    found = shutil.which("ollama")
    if found:
        return found
    if os.name == "nt":
        candidates = [
            Path(os.environ.get("LOCALAPPDATA", "")) / "Programs/Ollama/ollama.exe",
            Path(os.environ.get("ProgramFiles", "C:/Program Files")) / "Ollama/ollama.exe",
        ]
        for candidate in candidates:
            if candidate.is_file():
                return str(candidate)
    return None


def _find_poppler_tool(name: str) -> str | None:
    executable = f"{name}.exe" if os.name == "nt" else name
    candidates: list[Path] = []
    if os.name == "nt":
        packages_root = Path(os.environ.get("LOCALAPPDATA", "")) / "Microsoft/WinGet/Packages"
        if packages_root.is_dir():
            candidates.extend(
                packages_root.glob(
                    f"oschwartz10612.Poppler_*/poppler-*/Library/bin/{executable}"
                )
            )
        local_app_data = Path(os.environ.get("LOCALAPPDATA", ""))
        candidates.append(local_app_data / f"Programs/MiKTeX/miktex/bin/x64/{executable}")
    for candidate in candidates:
        if candidate.is_file():
            return str(candidate)
    found = shutil.which(name)
    return found if found and Path(found).exists() else None


def _read_pdf_pages(file_path: Path) -> tuple[list[str], str]:
    fallback = _read_pdf_with_pypdf(file_path)
    pdftotext = _find_poppler_tool("pdftotext")
    if not pdftotext:
        return fallback, "pypdf"
    with tempfile.TemporaryDirectory() as temporary:
        output = Path(temporary) / "layout.txt"
        completed = subprocess.run(
            [pdftotext, "-layout", "-enc", "UTF-8", str(file_path), str(output)],
            capture_output=True,
            text=True,
            timeout=180,
            check=False,
        )
        if completed.returncode != 0 or not output.is_file():
            return fallback, "pypdf"
        text = output.read_text(encoding="utf-8", errors="replace")
    extracted = text.replace("\r\n", "\n").split("\f")
    while extracted and not extracted[-1].strip():
        extracted.pop()
    if len(extracted) != len(fallback):
        return fallback, "pypdf"
    pages = [
        poppler_page.strip() if poppler_page.strip() else fallback[index]
        for index, poppler_page in enumerate(extracted)
    ]
    return pages, "Poppler pdftotext -layout"


def analyze_pdf(file_path: Path) -> dict[str, object]:
    """Inspect provenance, raster overlays, embedded images, vectors, and text."""
    reader = PdfReader(str(file_path))
    metadata = reader.metadata or {}
    page_count = len(reader.pages)
    first_page_size: list[float] | None = None
    if page_count:
        media_box = reader.pages[0].mediabox
        first_page_size = [float(media_box.width), float(media_box.height)]

    text_characters = 0
    text_pages = 0
    embedded_images = 0
    pages_with_images = 0
    full_page_raster_pages = 0
    vector_drawing_count = 0
    if fitz is not None:
        document = fitz.open(file_path)  # type: ignore[union-attr]
        try:
            for page_number in range(document.page_count):
                page = document.load_page(page_number)
                page_text = page.get_text("text") or ""
                text_characters += len(page_text.strip())
                if page_text.strip():
                    text_pages += 1
                images = page.get_images(full=True)
                embedded_images += len(images)
                if images:
                    pages_with_images += 1
                page_area = max(float(page.rect.width * page.rect.height), 1.0)
                has_full_page_raster = False
                for image in images:
                    xref = image[0]
                    try:
                        rectangles = page.get_image_rects(xref)
                    except Exception:  # noqa: BLE001
                        rectangles = []
                    for rectangle in rectangles:
                        if float(rectangle.width * rectangle.height) / page_area >= 0.85:
                            has_full_page_raster = True
                            break
                    if has_full_page_raster:
                        break
                if has_full_page_raster:
                    full_page_raster_pages += 1
                try:
                    vector_drawing_count += len(page.get_drawings())
                except Exception:  # noqa: BLE001
                    pass
        finally:
            document.close()

    pdfinfo_values: dict[str, str] = {}
    pdfinfo = _find_poppler_tool("pdfinfo")
    if pdfinfo:
        completed = subprocess.run(
            [pdfinfo, str(file_path)],
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=60,
            check=False,
        )
        if completed.returncode == 0:
            for line in completed.stdout.splitlines():
                if ":" in line:
                    key, value = line.split(":", 1)
                    pdfinfo_values[key.strip()] = value.strip()

    raster_ratio = full_page_raster_pages / page_count if page_count else 0.0
    text_ratio = text_pages / page_count if page_count else 0.0
    warnings: list[str] = []
    if raster_ratio >= 0.8 and text_ratio >= 0.8:
        warnings.append(
            "Most pages contain both a full-page raster and extractable text. This may be a scanned/OCR PDF or a previously generated overlay; compare producer/date/size with any alternate source PDF before transcribing."
        )
    elif raster_ratio >= 0.8 and text_ratio < 0.2:
        warnings.append("This is primarily a raster/scanned PDF; OCR or a vision model is recommended.")
    elif text_ratio < 0.2:
        warnings.append("Very little extractable text was found; inspect the source and consider OCR.")
    return {
        "file": str(file_path),
        "page_count": page_count,
        "page_size_points": first_page_size,
        "title": str(metadata.get("/Title") or pdfinfo_values.get("Title") or ""),
        "creator": str(metadata.get("/Creator") or pdfinfo_values.get("Creator") or ""),
        "producer": str(metadata.get("/Producer") or pdfinfo_values.get("Producer") or ""),
        "creation_date": str(
            metadata.get("/CreationDate") or pdfinfo_values.get("CreationDate") or ""
        ),
        "text_characters": text_characters,
        "text_pages": text_pages,
        "embedded_images": embedded_images,
        "pages_with_images": pages_with_images,
        "full_page_raster_pages": full_page_raster_pages,
        "vector_drawing_count": vector_drawing_count,
        "raster_page_ratio": raster_ratio,
        "text_page_ratio": text_ratio,
        "likely_full_page_raster_document": raster_ratio >= 0.8,
        "likely_raster_with_text_overlay": raster_ratio >= 0.8 and text_ratio >= 0.8,
        "warnings": warnings,
        "poppler_available": bool(pdfinfo),
    }


def _ollama_base_url(endpoint: str) -> str:
    parsed = urllib.parse.urlsplit(endpoint.strip() or DEFAULT_OLLAMA_ENDPOINT)
    if not parsed.scheme or not parsed.netloc:
        return "http://127.0.0.1:11434"
    return f"{parsed.scheme}://{parsed.netloc}"


def _version_tuple(version: str) -> tuple[int, ...]:
    numbers = re.findall(r"\d+", version.split("-", 1)[0])
    return tuple(int(number) for number in numbers[:4]) or (0,)


def check_for_updates(timeout: float = 8.0) -> dict[str, object]:
    """Check the public GitHub release feed without changing the installation."""
    request = urllib.request.Request(
        GITHUB_LATEST_RELEASE_API,
        headers={
            "Accept": "application/vnd.github+json",
            "User-Agent": f"book-to-latex/{APP_VERSION}",
        },
    )
    try:
        with urllib.request.urlopen(request, timeout=timeout) as response:
            payload = json.loads(response.read().decode("utf-8"))
        latest = str(payload.get("tag_name") or "").lstrip("v")
        assets = [
            {
                "name": str(asset.get("name") or ""),
                "url": str(asset.get("browser_download_url") or ""),
            }
            for asset in payload.get("assets", [])
            if isinstance(asset, dict)
        ]
        return {
            "success": bool(latest),
            "current_version": APP_VERSION,
            "latest_version": latest,
            "update_available": bool(latest) and _version_tuple(latest) > _version_tuple(APP_VERSION),
            "release_url": str(payload.get("html_url") or "https://github.com/kkmmee94/book-to-latex/releases/latest"),
            "assets": assets,
            "error": None,
        }
    except Exception as exc:  # noqa: BLE001
        return {
            "success": False,
            "current_version": APP_VERSION,
            "latest_version": "",
            "update_available": False,
            "release_url": "https://github.com/kkmmee94/book-to-latex/releases/latest",
            "assets": [],
            "error": str(exc),
        }


def ollama_connection_info(
    endpoint: str = DEFAULT_OLLAMA_ENDPOINT,
    timeout: float = 3.0,
) -> dict[str, object]:
    """Return Ollama availability and installed models without raising UI errors."""
    base_url = _ollama_base_url(endpoint)
    try:
        with urllib.request.urlopen(f"{base_url}/api/version", timeout=timeout) as response:
            version_payload = json.loads(response.read().decode("utf-8"))
        with urllib.request.urlopen(f"{base_url}/api/tags", timeout=timeout) as response:
            tags_payload = json.loads(response.read().decode("utf-8"))
        models = []
        for model in tags_payload.get("models", []):
            name = str(model.get("name") or model.get("model") or "").strip()
            if not name:
                continue
            details = model.get("details") or {}
            models.append(
                {
                    "name": name,
                    "size": int(model.get("size") or 0),
                    "parameter_size": str(details.get("parameter_size") or ""),
                    "quantization": str(details.get("quantization_level") or ""),
                }
            )
        models.sort(key=lambda item: str(item["name"]).lower())
        return {
            "available": True,
            "base_url": base_url,
            "version": str(version_payload.get("version") or ""),
            "models": models,
            "error": None,
        }
    except Exception as exc:  # noqa: BLE001
        return {
            "available": False,
            "base_url": base_url,
            "version": "",
            "models": [],
            "error": str(exc),
        }


def runtime_capabilities() -> dict[str, object]:
    """Return optional tool availability for friendly UI diagnostics."""
    tesseract_path = _find_tesseract()
    ocr_languages: list[str] = []
    if pytesseract is not None and tesseract_path:
        try:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
            ocr_languages = sorted(pytesseract.get_languages(config=""))
        except Exception:  # noqa: BLE001
            ocr_languages = []
    bundled_tessdata = _bundled_tessdata_dir()
    if bundled_tessdata:
        ocr_languages = sorted(
            set(ocr_languages)
            | {path.stem for path in bundled_tessdata.glob("*.traineddata")}
        )
    user_tessdata = _user_tessdata_dir()
    if user_tessdata.is_dir():
        ocr_languages = sorted(
            set(ocr_languages)
            | {path.stem for path in user_tessdata.glob("*.traineddata")}
        )
    return {
        "pdf_compiler": shutil.which("pdflatex"),
        "unicode_pdf_compiler": shutil.which("xelatex"),
        "pymupdf": fitz is not None,
        "pillow": _PILImage is not None,
        "pytesseract": pytesseract is not None,
        "tesseract": tesseract_path,
        "ocr_languages": ocr_languages,
        "libreoffice": _find_libreoffice(),
        "ollama": _find_ollama(),
        "poppler": {
            name: _find_poppler_tool(name)
            for name in ("pdfinfo", "pdftotext", "pdftoppm", "pdfimages")
        },
        "supported_extensions": sorted(SUPPORTED_INPUT_EXTENSIONS),
    }


def _ensure_ocr_dependencies() -> None:
    _ensure_fitz()
    _ensure_pillow()
    if pytesseract is None:
        raise RuntimeError(
            "pytesseract is required for OCR. Install with `pip install pytesseract` and Tesseract OCR."
        )
    tesseract_path = _find_tesseract()
    if not tesseract_path:
        raise RuntimeError(
            "Tesseract OCR is not installed or could not be found. Install Tesseract OCR, then restart the app."
        )
    pytesseract.pytesseract.tesseract_cmd = tesseract_path


def _render_pdf_page_image(
    pdf_doc: fitz.Document,
    page_no: int,
    dpi: int,
    color: bool,
    output_dir: Path,
) -> Path:
    _ensure_fitz()
    page = pdf_doc.load_page(page_no)
    matrix = fitz.Matrix(dpi / 72, dpi / 72)
    colorspace = fitz.csRGB if color else fitz.csGRAY
    pix = page.get_pixmap(matrix=matrix, colorspace=colorspace, alpha=False)
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"page_{page_no + 1:03d}.png"
    pix.save(str(output_path))
    return output_path


def _detect_language_from_text(text: str) -> str | None:
    sample = text[:100_000]
    counts = {
        "ara": sum("\u0600" <= char <= "\u06ff" for char in sample),
        "chi_sim": sum("\u4e00" <= char <= "\u9fff" for char in sample),
    }
    language, count = max(counts.items(), key=lambda item: item[1])
    return language if count >= 4 else None


def _detect_language_from_image(image: object) -> str | None:
    _ensure_ocr_dependencies()
    try:
        osd = pytesseract.image_to_osd(image)
    except Exception:  # noqa: BLE001
        return None
    match = re.search(r"Script:\s*([^\r\n]+)", osd, flags=re.I)
    script = match.group(1).strip().lower() if match else ""
    script_map = {
        "arabic": "ara",
        "han": "chi_sim",
        "han simplified": "chi_sim",
        "han traditional": "chi_tra",
        "latin": "eng",
    }
    return script_map.get(script)


def _extract_pdf_page_assets(
    pdf_doc: fitz.Document,
    page_no: int,
    output_dir: Path,
    asset_dir_name: str,
) -> list[tuple[Path, str]]:
    """Extract non-page-sized raster assets that LaTeX may need to retain."""
    page = pdf_doc.load_page(page_no)
    page_area = max(float(page.rect.width * page.rect.height), 1.0)
    assets: list[tuple[Path, str]] = []
    seen_xrefs: set[int] = set()
    for image_info in page.get_images(full=True):
        xref = int(image_info[0])
        if xref in seen_xrefs:
            continue
        seen_xrefs.add(xref)
        try:
            rectangles = page.get_image_rects(xref)
        except Exception:  # noqa: BLE001
            rectangles = []
        if not rectangles:
            continue
        if max(float(rect.width * rect.height) / page_area for rect in rectangles) >= 0.80:
            # This is usually a scan or existing full-page overlay, not a photo
            # that should be duplicated inside an editable reconstruction.
            continue
        try:
            extracted = pdf_doc.extract_image(xref)
        except Exception:  # noqa: BLE001
            continue
        width = int(extracted.get("width") or 0)
        height = int(extracted.get("height") or 0)
        image_bytes = extracted.get("image")
        extension = str(extracted.get("ext") or "png").lower()
        if not image_bytes or width < 32 or height < 32:
            continue
        if extension not in {"png", "jpg", "jpeg"}:
            try:
                pixmap = fitz.Pixmap(pdf_doc, xref)
                if pixmap.n - pixmap.alpha > 3:
                    pixmap = fitz.Pixmap(fitz.csRGB, pixmap)
                image_bytes = pixmap.tobytes("png")
                extension = "png"
            except Exception:  # noqa: BLE001
                continue
        output_dir.mkdir(parents=True, exist_ok=True)
        filename = f"page_{page_no + 1:03d}_asset_{xref}.{extension}"
        output_path = output_dir / filename
        if not output_path.is_file():
            output_path.write_bytes(bytes(image_bytes))
        reference = f"{asset_dir_name}/{filename}"
        assets.append((output_path, reference))
    return assets


def _save_image_frame_asset(
    file_path: Path,
    frame_number: int,
    page_number: int,
    output_dir: Path,
    asset_dir_name: str,
) -> tuple[Path, str]:
    """Save an image input in a LaTeX-compatible format for keep-photo mode."""
    frame = _load_image_frame(file_path, frame_number).convert("RGB")
    output_dir.mkdir(parents=True, exist_ok=True)
    filename = f"source_image_{page_number:03d}.png"
    output_path = output_dir / filename
    frame.save(output_path, format="PNG", optimize=True)
    return output_path, f"{asset_dir_name}/{filename}"


def _cleanup_unreferenced_assets(
    assets_dir: Path,
    all_assets: list[tuple[Path, str]],
    latex_body: str,
) -> list[str]:
    """Keep only assets that the generated LaTeX actually references."""
    kept: list[str] = []
    for path, reference in all_assets:
        if reference in latex_body:
            kept.append(reference)
        elif path.is_file():
            path.unlink()
    if assets_dir.is_dir() and not any(assets_dir.iterdir()):
        assets_dir.rmdir()
    return kept


def _ocr_pdf_page(
    pdf_doc: fitz.Document,
    page_no: int,
    dpi: int,
    lang: str,
    color: bool,
) -> str:
    _ensure_ocr_dependencies()
    ensure_ocr_language(lang or "eng")
    page = pdf_doc.load_page(page_no)
    matrix = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=matrix, alpha=False)
    assert _PILImage is not None
    img = _PILImage.frombytes("RGB", (pix.width, pix.height), pix.samples)
    if not color:
        img = img.convert("L")
    return pytesseract.image_to_string(
        img,
        lang=lang or "eng",
        config=_tesseract_config(lang or "eng"),
    )


def _image_frame_count(file_path: Path) -> int:
    _ensure_pillow()
    assert _PILImage is not None
    with _PILImage.open(file_path) as image:
        return max(int(getattr(image, "n_frames", 1)), 1)


def _load_image_frame(file_path: Path, frame_number: int) -> object:
    _ensure_pillow()
    assert _PILImage is not None
    with _PILImage.open(file_path) as image:
        image.seek(frame_number)
        frame = image.copy()
    if _PILImageOps is not None:
        frame = _PILImageOps.exif_transpose(frame)
    return frame


def _ocr_image_frame(file_path: Path, frame_number: int, lang: str) -> str:
    _ensure_ocr_dependencies()
    ensure_ocr_language(lang or "eng")
    frame = _load_image_frame(file_path, frame_number)
    return pytesseract.image_to_string(
        frame,
        lang=lang or "eng",
        config=_tesseract_config(lang or "eng"),
    )


def _render_image_frame(
    file_path: Path,
    frame_number: int,
    page_number: int,
    color: bool,
    output_dir: Path,
) -> Path:
    frame = _load_image_frame(file_path, frame_number)
    frame = frame.convert("RGB" if color else "L")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"page_{page_number:03d}.png"
    frame.save(output_path, format="PNG", optimize=True)
    return output_path


def _query_llm(
    *,
    page_text: str,
    page_no: int,
    source_format: str,
    page_image_path: Path | None,
    style_guide: str,
    redraw_graphs: bool,
    model: str,
    endpoint: str,
    api_key: str,
    max_tokens: int,
    temperature: float,
    timeout: float,
    retries: int,
    backoff: float,
    strict_mode: bool,
    match_mode: str,
    log_callback: LogCallback | None,
    document_language: str = "eng",
    preserve_layout: bool = False,
    preserve_color: bool = False,
    page_flow: str = PAGE_FLOW_SOURCE,
    photo_handling: str = PHOTO_KEEP,
    page_asset_references: list[str] | None = None,
    visual_inventory_summary: str = "",
    semantic_asset_references: list[str] | None = None,
    retained_visual_references: list[str] | None = None,
    force_semantic_retry: bool = False,
) -> str:
    system_prompt = [
        "You are a publication-quality LaTeX converter.",
        "Convert the extracted content into valid LaTeX body content only.",
        "Do not emit a document class, package imports, document environment, or Markdown fences. Use normal LaTeX sectioning commands when the source has headings.",
        "Preserve meaning and symbols exactly and do not insert explanations.",
        "Preserve headings, lists, slide structure, and tab-separated table rows when they are present.",
        DEFAULT_LATEX_STYLE_GUIDE,
        "Do not output chain-of-thought or hidden analysis.",
    ]

    language_name = DOCUMENT_LANGUAGES.get(document_language, document_language)
    system_prompt.append(
        f"The document language is {language_name}. Preserve the original language exactly; do not translate or transliterate it."
    )
    if document_language == "ara":
        system_prompt.append(
            "Preserve Arabic right-to-left reading order and Arabic punctuation. Return Unicode Arabic suitable for XeLaTeX with polyglossia."
        )

    if strict_mode:
        system_prompt.append(
            "Strict mode active: keep formula/math symbols valid for LaTeX and normalize Unicode math punctuation conservatively."
        )

    if match_mode == MATCH_MODE_EXACT:
        system_prompt.append("Avoid paraphrasing; preserve token order where possible.")
    else:
        system_prompt.append("Small formatting improvements are acceptable if text stays identical.")

    if style_guide.strip():
        system_prompt.append(f"Project-specific style guide:\n{style_guide.strip()[:20000]}")
    if page_image_path is not None:
        system_prompt.append(
            "A rendered source-page image is attached for analysis only. Read its full visual structure and mathematics, and cross-check every digit and word against the extracted text. Reconstruct stacked fractions, limits, superscripts, subscripts, matrices, cases, aligned derivations, tables, charts, timelines, flowcharts, geometric figures, technical diagrams, and other visuals that LaTeX can represent. Never emit an includegraphics command for the temporary page image and never invent an image filename."
        )
    if visual_inventory_summary:
        system_prompt.append(
            "A separate visual inventory identified the following page elements:\n"
            f"{visual_inventory_summary}"
        )
    if redraw_graphs:
        system_prompt.append(
            "Semantic visual reconstruction is required. When a graph, chart, table, timeline, flowchart, equation, geometric figure, or technical diagram is visible, recreate it as editable LaTeX using tabular, TikZ, pgfplots, or standard mathematical environments. Use only values, labels, shapes, and relationships that are visible or extracted. Never replace a reconstructible visual with a blank box, prose-only summary, or invented values."
        )

    if preserve_layout and page_flow == PAGE_FLOW_SOURCE:
        system_prompt.append(
            "Reproduce the source page structure closely, including meaningful title bands, columns, colored boxes, logos, headers, footers, and top/bottom content. Do not emit a page break; the app controls page boundaries."
        )
    elif page_flow == PAGE_FLOW_COMPACT:
        system_prompt.append(
            "Create a compact continuous document. Preserve content order and hierarchy, but omit repeated running headers, repeated footers, source page numbers, and empty spacing. Do not emit a page break; let LaTeX flow content naturally."
        )
    else:
        system_prompt.append(
            "Keep the source unit self-contained without emitting a page break; the app controls page boundaries."
        )

    system_prompt.append(
        "Preserve visible colors accurately with xcolor/TikZ styling."
        if preserve_color
        else "Render reconstructed text, rules, tables, charts, and diagrams in black and white."
    )

    asset_references = page_asset_references or []
    if asset_references:
        exact_assets = "\n".join(f"- {reference}" for reference in asset_references)
        system_prompt.append(
            "These are the only persistent raster assets available for this source page:\n"
            f"{exact_assets}\n"
            "You may reference them with \\includegraphics using the exact path shown. Never change, shorten, or invent a path. Logos, signatures, and meaningful artwork that cannot be represented faithfully in LaTeX may use these assets."
        )
    semantic_references = semantic_asset_references or []
    if semantic_references:
        system_prompt.append(
            "The following raster references contain graphs, charts, tables, infographics, or technical diagrams and MUST NOT be used with includegraphics. Recreate their visible information semantically in LaTeX/TikZ/pgfplots/tabular instead:\n"
            + "\n".join(f"- {reference}" for reference in semantic_references)
        )
    retained_references = retained_visual_references or []
    if retained_references:
        system_prompt.append(
            "The following semantic visuals cannot be reconstructed exactly because their underlying values or geometry are not fully recoverable. Keep them visible with includegraphics using the exact supplied path; do not invent an approximate redraw:\n"
            + "\n".join(f"- {reference}" for reference in retained_references)
        )
    if force_semantic_retry:
        system_prompt.append(
            "CORRECTION REQUIRED: a previous draft omitted or raster-embedded reconstructible visual content. This response must visibly recreate every inventoried semantic visual with actual LaTeX code."
        )
    if photo_handling == PHOTO_DESCRIBE:
        system_prompt.append(
            "For natural photographs of people, animals, places, or real-world objects, do not include the photograph. Replace each significant photograph at its approximate source position with a concise, objective italic description beginning with 'Photo description:'. Do not describe graphs, tables, equations, or diagrams; reconstruct those semantically instead."
        )
    else:
        system_prompt.append(
            "Keep significant natural photographs of people, animals, places, or real-world objects using their exact persistent asset path when one is available. Do not convert photographs into TikZ. Reconstruct graphs, tables, equations, and technical diagrams semantically instead of embedding their raster image."
        )

    user_prompt = (
        f"Source format: {source_format}\n"
        f"Content unit: {page_no + 1}\n\n"
        f"Page usage: {page_flow}\n"
        f"Natural-photo policy: {photo_handling}\n\n"
        "Extracted source content:\n"
        f"{page_text}\n\n"
        "Return only the LaTeX snippet for this page. /no_think"
    )

    headers = {"Content-Type": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    user_content: object = user_prompt
    encoded_image = ""
    if page_image_path is not None:
        encoded_image = base64.b64encode(page_image_path.read_bytes()).decode("ascii")
        user_content = [
            {"type": "text", "text": user_prompt},
            {
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{encoded_image}"},
            },
        ]

    endpoint_path = urllib.parse.urlsplit(endpoint).path.rstrip("/")
    ollama_native = endpoint_path == "/api/chat"
    if ollama_native:
        user_message: dict[str, object] = {"role": "user", "content": user_prompt}
        if encoded_image:
            user_message["images"] = [encoded_image]
        payload = {
            "model": model,
            "messages": [
                {"role": "system", "content": "\n".join(system_prompt)},
                user_message,
            ],
            "options": {"temperature": temperature, "num_predict": max_tokens},
            "think": False,
            "stream": False,
        }
    else:
        payload = {
            "model": model,
            "messages": [
                {"role": "system", "content": "\n".join(system_prompt)},
                {"role": "user", "content": user_content},
            ],
            "temperature": temperature,
            "max_tokens": max_tokens,
            "stream": False,
        }

    attempts = 0
    while True:
        attempts += 1
        try:
            _log(log_callback, f"LLM request for page {page_no + 1} (attempt {attempts})")
            request_payload = json.dumps(payload).encode("utf-8")
            request = urllib.request.Request(
                endpoint,
                data=request_payload,
                headers=headers,
                method="POST",
            )
            with urllib.request.urlopen(request, timeout=timeout) as response:
                payload_json = json.loads(response.read().decode("utf-8"))
            if ollama_native:
                content = payload_json.get("message", {}).get("content")
            else:
                content = (
                    payload_json.get("choices", [{}])[0]
                    .get("message", {})
                    .get("content")
                )
            if content is None:
                expected = "message.content" if ollama_native else "choices[0].message.content"
                raise RuntimeError(f"LLM response did not include {expected}")
            return str(content).strip()
        except Exception as exc:  # noqa: BLE001
            if attempts > retries:
                raise RuntimeError(f"LLM request failed for page {page_no + 1}: {exc}") from exc
            _log(log_callback, f"Retrying page {page_no + 1} after error: {exc}")
            time.sleep(backoff * attempts)


def _query_visual_inventory(
    *,
    page_image_path: Path,
    page_no: int,
    asset_references: list[str],
    model: str,
    endpoint: str,
    api_key: str,
    timeout: float,
    retries: int,
    backoff: float,
    log_callback: LogCallback | None,
) -> list[dict[str, str]]:
    """Classify extracted raster assets before deciding redraw/keep/describe."""
    system_prompt = (
        "You are a visual document inventory classifier. Match each supplied asset reference to what it represents on the attached page. "
        "Return JSON only with this exact shape: "
        '{"assets":[{"reference":"exact supplied path","kind":"graph|chart|table|diagram|infographic|equation|photo|logo|screenshot|artwork|other","description":"objective description","reconstructable":"yes|no"}]}. '
        "Use photo only for real people, animals, places, or physical objects. Graphs, charts, tables, infographics, equations and technical diagrams are semantic visuals even when stored as photographs or scans. Set reconstructable to yes only when all values, labels, and required geometry are visibly readable without guessing; use no for dense scientific plots, pictorial infographics, or any visual whose data cannot be recovered exactly. Include every supplied reference exactly once."
    )
    user_prompt = (
        f"Source page: {page_no + 1}\n"
        "Asset references:\n"
        + "\n".join(f"- {reference}" for reference in asset_references)
    )
    encoded_image = base64.b64encode(page_image_path.read_bytes()).decode("ascii")
    endpoint_path = urllib.parse.urlsplit(endpoint).path.rstrip("/")
    ollama_native = endpoint_path == "/api/chat"
    headers = {"Content-Type": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"
    if ollama_native:
        payload = {
            "model": model,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt, "images": [encoded_image]},
            ],
            "format": "json",
            "options": {"temperature": 0, "num_predict": 900},
            "think": False,
            "stream": False,
        }
    else:
        payload = {
            "model": model,
            "messages": [
                {"role": "system", "content": system_prompt},
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": user_prompt},
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/png;base64,{encoded_image}"},
                        },
                    ],
                },
            ],
            "temperature": 0,
            "max_tokens": 900,
            "stream": False,
        }

    attempts = 0
    while True:
        attempts += 1
        try:
            _log(log_callback, f"Classifying visuals on page {page_no + 1}")
            request = urllib.request.Request(
                endpoint,
                data=json.dumps(payload).encode("utf-8"),
                headers=headers,
                method="POST",
            )
            with urllib.request.urlopen(request, timeout=timeout) as response:
                response_json = json.loads(response.read().decode("utf-8"))
            if ollama_native:
                content = str(response_json.get("message", {}).get("content") or "")
            else:
                content = str(
                    response_json.get("choices", [{}])[0]
                    .get("message", {})
                    .get("content")
                    or ""
                )
            cleaned = re.sub(r"^```(?:json)?\s*|\s*```$", "", content.strip(), flags=re.I)
            parsed = json.loads(cleaned)
            raw_assets = parsed.get("assets", []) if isinstance(parsed, dict) else []
            by_reference = {
                str(item.get("reference") or ""): item
                for item in raw_assets
                if isinstance(item, dict)
            }
            inventory: list[dict[str, str]] = []
            for reference in asset_references:
                item = by_reference.get(reference) or {}
                inventory.append(
                    {
                        "reference": reference,
                        "kind": str(item.get("kind") or "other").lower(),
                        "description": str(item.get("description") or "Visual from the source page").strip(),
                        "reconstructable": str(item.get("reconstructable") or "no").lower(),
                    }
                )
            return inventory
        except Exception as exc:  # noqa: BLE001
            if attempts > retries:
                _log(
                    log_callback,
                    f"Page {page_no + 1}: visual classification failed; keeping unclassified assets available ({exc})",
                )
                return [
                    {
                        "reference": reference,
                        "kind": "other",
                        "description": "Visual from the source page",
                        "reconstructable": "no",
                    }
                    for reference in asset_references
                ]
            time.sleep(backoff * attempts)


def _build_header(
    *,
    preserve_layout: bool,
    preserve_color: bool,
    preserve_graphs: bool,
    exact_visual_mode: bool,
    source_is_pdf: bool,
    document_language: str,
    page_size: str = PAGE_SIZE_A4,
    page_flow: str = PAGE_FLOW_SOURCE,
    source_page_size_points: list[float] | None = None,
    has_assets: bool = False,
    redraw_graphs: bool = False,
    images_dir: str | None = None,
) -> str:
    lines: list[str] = []
    paper_option = {
        PAGE_SIZE_A4: "a4paper",
        PAGE_SIZE_LETTER: "letterpaper",
    }.get(page_size)
    if page_size == PAGE_SIZE_SOURCE and not source_page_size_points:
        paper_option = "a4paper"
    class_options = "11pt" + (f",{paper_option}" if paper_option else "")
    class_name = (
        "article"
        if source_is_pdf and preserve_layout and not exact_visual_mode
        else "book"
    )
    lines.append(f"\\documentclass[{class_options}]{{{class_name}}}")
    if document_language == "ara":
        lines.append("\\usepackage{fontspec}")
        lines.append("\\usepackage{polyglossia}")
        lines.append("\\setdefaultlanguage{arabic}")
        lines.append("\\setotherlanguage{english}")
        lines.append("\\setmainfont{Amiri}")
        lines.append("\\newfontfamily\\arabicfont[Script=Arabic]{Amiri}")
    elif document_language in {"chi_sim", "chi_tra"}:
        lines.append("\\usepackage{fontspec}")
        lines.append("\\usepackage[UTF8,scheme=plain,fontset=none]{ctex}")
        lines.append("\\setCJKmainfont{FandolSong-Regular}")
    else:
        lines.append("\\usepackage[utf8]{inputenc}")
        lines.append("\\usepackage[T1]{fontenc}")
        lines.append("\\usepackage{lmodern}")
    lines.append("\\usepackage{microtype}")
    lines.append("\\usepackage{textcomp}")
    lines.append("\\usepackage{amsmath,amssymb,mathtools}")
    lines.append("\\usepackage{booktabs,array}")
    lines.append("\\usepackage{enumitem}")
    lines.append("\\usepackage{listings}")
    lines.append("\\usepackage[hidelinks]{hyperref}")
    lines.append("\\setlist[enumerate,1]{label=(\\alph*)}")
    lines.append("\\setlist[enumerate,2]{label=\\roman*.}")
    lines.append("\\setlist[enumerate,3]{label=\\Alph*.}")
    lines.append("\\lstset{basicstyle=\\ttfamily\\small,breaklines=true,columns=fullflexible}")
    lines.append("\\lstnewenvironment{Rcode}{}{}")
    lines.append("\\providecommand{\\Prob}[2]{\\section*{Problem #1: #2}\\label{prob:#1}\\addcontentsline{toc}{section}{Problem #1: #2}\\hfill\\hyperref[sol:#1]{Go to solution}}")
    lines.append("\\providecommand{\\Sol}[2]{\\section*{Solution #1: #2}\\label{sol:#1}\\addcontentsline{toc}{section}{Solution #1: #2}\\hfill\\hyperref[prob:#1]{Go to problem}}")

    if redraw_graphs:
        lines.append("\\usepackage{tikz}")
        lines.append("\\usepackage{pgfplots}")
        lines.append("\\usepgfplotslibrary{statistics}")
        lines.append("\\pgfplotsset{compat=1.18}")

    if exact_visual_mode and source_is_pdf:
        lines.append("\\usepackage{pdfpages}")
    elif preserve_graphs or exact_visual_mode or has_assets:
        lines.append("\\usepackage{graphicx}")
        lines.append("\\usepackage{float}")
    if preserve_layout and page_flow == PAGE_FLOW_SOURCE and not exact_visual_mode:
        lines.append("\\usepackage{adjustbox}")

    if preserve_color or redraw_graphs:
        lines.append("\\usepackage[dvipsnames,table]{xcolor}")
    elif preserve_graphs:
        lines.append("\\usepackage{xcolor}")
        lines.append("\\AtBeginDocument{\\color{black}}")

    source_is_landscape = bool(
        source_page_size_points
        and len(source_page_size_points) >= 2
        and source_page_size_points[0] > source_page_size_points[1]
    )
    if preserve_layout and source_is_landscape and document_language != "ara":
        lines.append("\\usepackage{titlesec}")
        lines.append("\\renewcommand{\\familydefault}{\\sfdefault}")
        lines.append("\\titleformat{\\section}{\\Large\\sffamily\\mdseries}{}{0pt}{}")
        lines.append("\\titlespacing*{\\section}{0pt}{0pt}{0.45em}")

    if page_size == PAGE_SIZE_SOURCE and source_page_size_points:
        paper_width, paper_height = source_page_size_points
        margin = "0.25in" if preserve_layout else "0.7in"
        lines.append(
            "\\usepackage["
            f"paperwidth={paper_width:.3f}bp,paperheight={paper_height:.3f}bp,margin={margin}"
            "]{geometry}"
        )
    elif preserve_layout:
        lines.append("\\usepackage[margin=0.65in]{geometry}")
    else:
        lines.append("\\usepackage[margin=1.4in]{geometry}")
    if preserve_layout:
        lines.append("\\setlength{\\parskip}{0.3em}")
        lines.append("\\setlength{\\parindent}{0pt}")
    else:
        lines.append("\\setlength{\\parskip}{0.8em}")
        lines.append("\\setlength{\\parindent}{0pt}")

    if images_dir:
        safe_dir = images_dir.replace("\\", "/")
        lines.append("\\graphicspath{{" + safe_dir + "/}}")

    lines.append("")
    lines.append("\\begin{document}")
    if preserve_layout and page_flow == PAGE_FLOW_SOURCE:
        lines.append("\\pagestyle{empty}")
    if not source_is_pdf and not exact_visual_mode:
        lines.append("\\chapter*{Converted manuscript}")

    lines.append("")
    return "\n".join(lines)


def _build_footer() -> str:
    return "\n\\end{document}"


def _constrain_source_page_layout(latex: str, preserve_layout: bool, page_flow: str) -> str:
    """Prevent reconstructed source pages from overflowing onto extra pages."""
    if not preserve_layout or page_flow != PAGE_FLOW_SOURCE:
        return latex

    graphics_pattern = re.compile(
        r"\\includegraphics(?:\[(?P<options>[^\]]*)\])?\{(?P<reference>[^{}]+)\}"
    )

    def constrain_graphic(match: re.Match[str]) -> str:
        options = [option.strip() for option in (match.group("options") or "").split(",") if option.strip()]
        if not any(option.startswith("height=") for option in options):
            options.append(r"height=0.68\textheight")
        if "keepaspectratio" not in options:
            options.append("keepaspectratio")
        if not any(option.startswith("width=") for option in options):
            options.insert(0, r"width=\linewidth")
        return f"\\includegraphics[{','.join(options)}]{{{match.group('reference')}}}"

    constrained = graphics_pattern.sub(constrain_graphic, latex)
    constrained = re.sub(
        r"\\vspace\*?\{(?:[1-9]\d*(?:\.\d+)?|0\.[5-9]\d*)\s*(?:cm|in|mm)\}",
        r"\\vspace{0.35em}",
        constrained,
    )
    return constrained


def _page_to_latex(
    page_no: int,
    page_text: str,
    source_format: str,
    style_guide: str,
    vision_mode: bool,
    redraw_graphs: bool,
    match_mode: str,
    match_error_percent: float,
    model: str,
    endpoint: str,
    api_key: str,
    max_tokens: int,
    temperature: float,
    timeout: float,
    retries: int,
    backoff: float,
    strict_mode: bool,
    no_llm: bool,
    preserve_graphs: bool,
    preserve_layout: bool,
    image_path: Path | None,
    image_reference: str | None,
    exact_visual_mode: bool,
    log_callback: LogCallback | None,
    document_language: str = "eng",
    exact_source_reference: str | None = None,
    page_size: str = PAGE_SIZE_SOURCE,
    page_flow: str = PAGE_FLOW_SOURCE,
    preserve_color: bool = False,
    photo_handling: str = PHOTO_KEEP,
    page_asset_references: list[str] | None = None,
) -> PageOutput:
    source = (page_text or "").strip()
    warnings: list[str] = []
    used_fallback = False
    image_ref = (image_reference or (image_path.name if image_path else "")).replace("\\", "/")

    if exact_visual_mode and exact_source_reference:
        source_reference = exact_source_reference.replace("\\", "/")
        if source_format.startswith("PDF"):
            fitpaper = "true" if page_size == PAGE_SIZE_SOURCE else "false"
            snippet = (
                f"\\includepdf[pages={{{page_no + 1}}},fitpaper={fitpaper},pagecommand={{}}]"
                f"{{\\detokenize{{{source_reference}}}}}"
            )
        else:
            snippet = (
                "\\begin{figure}[H]\n"
                "\\centering\n"
                "\\includegraphics[width=\\textwidth,height=0.96\\textheight,keepaspectratio]"
                f"{{\\detokenize{{{source_reference}}}}}\n"
                "\\end{figure}"
            )
        return PageOutput(
            page_no=page_no + 1,
            source_text=source,
            llm_text=snippet,
            latex_body=snippet,
            match_ratio=1.0,
            uncertain=False,
            image_path=None,
            warnings=None,
            used_fallback=False,
        )

    if exact_visual_mode and image_path is not None:
        snippet = (
            "\\begin{figure}[H]\n"
            "\\centering\n"
            f"\\includegraphics[width=\\textwidth,height=0.96\\textheight,keepaspectratio]{{{image_ref}}}\n"
            "\\end{figure}"
        )
        return PageOutput(
            page_no=page_no + 1,
            source_text=source,
            llm_text=snippet,
            latex_body=snippet,
            match_ratio=1.0,
            uncertain=False,
            image_path=image_path,
            warnings=None,
            used_fallback=False,
        )

    if no_llm:
        llm = _fallback_latex(
            source,
            preserve_layout=preserve_layout,
            strict_mode=True,
        )
    else:
        try:
            asset_references = page_asset_references or []
            inventory: list[dict[str, str]] = []
            if redraw_graphs and image_path is not None and asset_references:
                inventory = _query_visual_inventory(
                    page_image_path=image_path,
                    page_no=page_no,
                    asset_references=asset_references,
                    model=model,
                    endpoint=endpoint,
                    api_key=api_key,
                    timeout=timeout,
                    retries=retries,
                    backoff=backoff,
                    log_callback=log_callback,
                )
            semantic_kinds = {
                "graph",
                "chart",
                "table",
                "diagram",
                "infographic",
                "equation",
            }
            semantic_items = [item for item in inventory if item["kind"] in semantic_kinds]
            semantic_references = [
                item["reference"]
                for item in semantic_items
                if item.get("reconstructable") == "yes"
            ]
            retained_visual_items = [
                item
                for item in semantic_items
                if item.get("reconstructable") != "yes"
            ]
            retained_visual_references = [
                item["reference"] for item in retained_visual_items
            ]
            if retained_visual_items:
                warnings.append(
                    "One or more source visuals were kept as images because exact semantic reconstruction would require guessing missing data"
                )
            photo_items = [item for item in inventory if item["kind"] == "photo"]
            always_keep_kinds = {"logo", "screenshot", "artwork"}
            required_keep_items = [
                item
                for item in inventory
                if item["kind"] in always_keep_kinds
                or (item["kind"] == "photo" and photo_handling == PHOTO_KEEP)
                or item in retained_visual_items
            ]
            allowed_asset_references = [
                reference
                for reference in asset_references
                if reference not in semantic_references
                and not (
                    photo_handling == PHOTO_DESCRIBE
                    and any(
                        item["reference"] == reference and item["kind"] == "photo"
                        for item in inventory
                    )
                )
            ]
            inventory_summary = "\n".join(
                f"- {item['kind']} (exactly reconstructable: {item.get('reconstructable', 'no')}): {item['description']} (asset: {item['reference']})"
                for item in inventory
            )

            def request_latex(force_retry: bool = False) -> str:
                return _query_llm(
                    page_text=source,
                    page_no=page_no,
                    source_format=source_format,
                    page_image_path=image_path if vision_mode else None,
                    style_guide=style_guide,
                    redraw_graphs=redraw_graphs,
                    model=model,
                    endpoint=endpoint,
                    api_key=api_key,
                    max_tokens=max_tokens,
                    temperature=temperature,
                    timeout=timeout,
                    retries=retries,
                    backoff=backoff,
                    strict_mode=strict_mode,
                    match_mode=match_mode,
                    log_callback=log_callback,
                    document_language=document_language,
                    preserve_layout=preserve_layout,
                    preserve_color=preserve_color,
                    page_flow=page_flow,
                    photo_handling=photo_handling,
                    page_asset_references=allowed_asset_references,
                    visual_inventory_summary=inventory_summary,
                    semantic_asset_references=semantic_references,
                    retained_visual_references=retained_visual_references,
                    force_semantic_retry=force_retry,
                )

            llm_raw = request_latex()
            semantic_markers = (
                r"\begin{tikzpicture}",
                r"\begin{axis}",
                r"\begin{tabular}",
                r"\begin{longtable}",
                r"\begin{array}",
                r"\draw",
                r"\addplot",
            )
            ungrounded_markers = (
                "rand",
                "placeholder",
                "dummy data",
                "approximate visual",
                "approximation of",
                "schematic only",
                "...",
                r"\dots",
            )
            semantic_violation = bool(semantic_references) and (
                any(reference in llm_raw for reference in semantic_references)
                or not any(marker in llm_raw for marker in semantic_markers)
                or any(marker in llm_raw.lower() for marker in ungrounded_markers)
            )
            if semantic_violation:
                _log(
                    log_callback,
                    f"Page {page_no + 1}: retrying because a reconstructible visual was omitted or embedded as a raster",
                )
                llm_raw = request_latex(force_retry=True)
            llm = _clean_latex_output(llm_raw, strict_mode=strict_mode)
            llm, llm_issues = _post_validate_latex(llm)
            warnings.extend(llm_issues)
            fatal_validation = any(
                issue.startswith("Unbalanced braces") or issue.startswith("Unexpected")
                for issue in llm_issues
            )
            if not llm or fatal_validation:
                raise ValueError("AI returned invalid LaTeX")

            semantic_still_missing = bool(semantic_references) and (
                any(reference in llm for reference in semantic_references)
                or not any(marker in llm for marker in semantic_markers)
                or any(marker in llm.lower() for marker in ungrounded_markers)
            )
            if semantic_still_missing:
                warnings.append(
                    "AI could not recreate one or more semantic visuals; the source visual was retained so it remains visible"
                )
                figure_pattern = re.compile(
                    r"\\begin\{figure\}(?:\[[^\]]*\])?.*?\\end\{figure\}",
                    re.DOTALL,
                )

                def remove_ungrounded_figure(match: re.Match[str]) -> str:
                    figure_text = match.group(0)
                    if any(marker in figure_text for marker in semantic_markers):
                        return ""
                    if any(reference in figure_text for reference in semantic_references):
                        return ""
                    return figure_text

                llm = figure_pattern.sub(remove_ungrounded_figure, llm)
                for item in inventory:
                    if item["reference"] not in semantic_references:
                        continue
                    description = _escape_latex(item["description"])
                    llm += (
                        "\n\\begin{figure}[H]\n\\centering\n"
                        f"\\includegraphics[width=.9\\linewidth]{{{item['reference']}}}\n"
                        f"\\caption{{{description}}}\n\\end{{figure}}"
                    )

            for item in required_keep_items:
                if item["reference"] in llm:
                    continue
                description = _escape_latex(item["description"])
                llm += (
                    "\n\\begin{figure}[H]\n\\centering\n"
                    f"\\includegraphics[width=.9\\linewidth]{{{item['reference']}}}\n"
                    f"\\caption{{{description}}}\n\\end{{figure}}"
                )
            if photo_handling == PHOTO_DESCRIBE:
                for item in photo_items:
                    description = _escape_latex(item["description"])
                    if description in llm:
                        continue
                    llm += (
                        "\n\\begin{quote}\\textit{Photo description: "
                        f"{description}"
                        "}\\end{quote}"
                    )
        except Exception as exc:  # noqa: BLE001
            used_fallback = True
            warnings.append(f"AI conversion failed; safe local conversion used instead ({exc})")
            _log(log_callback, f"Page {page_no + 1}: AI unavailable, using safe local conversion")
            llm = _fallback_latex(
                source,
                preserve_layout=preserve_layout,
                strict_mode=True,
            )

    llm = _constrain_source_page_layout(llm, preserve_layout, page_flow)
    page_body = llm
    if preserve_graphs and image_path is not None:
        figure = (
            "\\begin{figure}[H]\n"
            "\\centering\n"
            f"\\includegraphics[width=\\linewidth]{{{image_ref}}}\n"
            "\\end{figure}"
        )
        page_body = f"{figure}\n{page_body}"
    if preserve_layout and page_flow == PAGE_FLOW_SOURCE:
        page_body = (
            "\\begin{adjustbox}{max width=\\textwidth,max totalheight=0.94\\textheight,center}\n"
            "\\begin{minipage}{\\textwidth}\n"
            f"{page_body}\n"
            "\\end{minipage}\n"
            "\\end{adjustbox}"
        )

    source_plain = _normalize_plain(source)
    generated_plain = _latex_to_plain(page_body)
    ratio = _match_similarity(source_plain, generated_plain)
    numeric_comparison = _numeric_comparison(source_plain, generated_plain)
    missing_numbers = numeric_comparison["missing_numbers"]
    extra_numbers = numeric_comparison["extra_numbers"]
    number_match = bool(numeric_comparison["numbers_match"])

    threshold = 1.0 if match_mode == MATCH_MODE_EXACT else max(0.0, 1.0 - (match_error_percent / 100.0))
    uncertain = ratio < threshold or used_fallback or not number_match
    if uncertain:
        if ratio < threshold:
            warnings.append(f"Text similarity {ratio:.2%} is below the {threshold:.2%} target")
        if not number_match:
            warnings.append(
                "Numeric values differ"
                f" (missing: {', '.join(missing_numbers) or 'none'};"
                f" extra: {', '.join(extra_numbers) or 'none'})"
            )

    return PageOutput(
        page_no=page_no + 1,
        source_text=source,
        llm_text=llm,
        latex_body=page_body,
        match_ratio=ratio,
        uncertain=uncertain,
        image_path=image_path,
        warnings="; ".join(warnings) if warnings else None,
        used_fallback=used_fallback,
        number_match=number_match,
        missing_numbers=", ".join(missing_numbers) if not number_match else "",
        extra_numbers=", ".join(extra_numbers) if not number_match else "",
    )


def _build_review_html(
    review_dir: Path,
    output_path: Path,
    page_outputs: list[PageOutput],
    compilation: dict[str, object] | None,
    verification: dict[str, object],
    pdf_analysis: dict[str, object] | None,
) -> str:
    uncertain_count = sum(1 for page in page_outputs if page.uncertain)
    compile_message = "PDF creation was not requested."
    compile_class = "muted"
    if compilation:
        if compilation.get("success"):
            compile_message = f"PDF created successfully: {compilation.get('pdf_path')}"
            compile_class = "good"
        else:
            compile_message = f"PDF was not created: {compilation.get('message', 'Unknown error')}"
            compile_class = "warning"

    pdf_diagnostic = ""
    if pdf_analysis:
        pdf_diagnostic = (
            "<div>Source PDF: "
            f"<strong>{pdf_analysis.get('page_count', 0)}</strong> pages · "
            f"{pdf_analysis.get('full_page_raster_pages', 0)} full-page rasters · "
            f"{pdf_analysis.get('vector_drawing_count', 0)} vector objects"
            "</div>"
        )

    cards: list[str] = []
    for page in page_outputs:
        state = "review" if page.uncertain else "good"
        label = "Needs review" if page.uncertain else "Passed"
        warning = html.escape(page.warnings or "No warnings")
        image_block = ""
        if page.image_path and page.image_path.is_file():
            encoded_source_page = base64.b64encode(page.image_path.read_bytes()).decode("ascii")
            image_block = (
                "<div><h3>Rendered source page</h3>"
                f'<img class="source-page" src="data:image/png;base64,{encoded_source_page}" alt="Rendered source page {page.page_no}"></div>'
            )
        cards.append(
            f"""
            <section class="page-card {state}" data-review="{'yes' if page.uncertain else 'no'}">
              <div class="page-heading">
                <h2>Page {page.page_no}</h2>
                <span class="badge {state}">{label}</span>
                <span class="ratio">Text match: {page.match_ratio:.2%}</span>
                <span class="ratio">Numbers: {'matched' if page.number_match else 'different'}</span>
              </div>
              <p class="warning-text">{warning}</p>
              <div class="columns">
                {image_block}
                <div><h3>Original text</h3><pre>{html.escape(page.source_text)}</pre></div>
                <div><h3>Generated LaTeX</h3><pre>{html.escape(page.latex_body)}</pre></div>
              </div>
            </section>
            """
        )

    return f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Book conversion review</title>
<style>
  :root {{ color-scheme: light; font-family: Inter, Segoe UI, Arial, sans-serif; }}
  body {{ margin: 0; background: #f4f6fa; color: #172033; }}
  header {{ padding: 28px max(24px, 5vw); color: white; background: linear-gradient(120deg,#243b6b,#4169a8); }}
  header h1 {{ margin: 0 0 8px; }} header p {{ margin: 5px 0; }}
  main {{ max-width: 1400px; margin: 0 auto; padding: 24px; }}
  .summary {{ display: flex; gap: 16px; flex-wrap: wrap; margin-bottom: 20px; }}
  .summary > div {{ background: white; border-radius: 10px; padding: 14px 18px; box-shadow: 0 2px 10px #17203312; }}
  button {{ border: 0; border-radius: 8px; padding: 10px 14px; background: #243b6b; color: white; cursor: pointer; }}
  .page-card {{ background: white; border-left: 6px solid #2b8a5a; border-radius: 10px; margin: 18px 0; padding: 18px; box-shadow: 0 2px 10px #17203312; }}
  .page-card.review {{ border-left-color: #d97706; }}
  .page-heading {{ display: flex; align-items: center; gap: 12px; flex-wrap: wrap; }}
  .page-heading h2 {{ margin: 0; }} .ratio {{ color: #526079; }}
  .badge {{ border-radius: 999px; padding: 4px 9px; font-size: 0.84rem; background: #dff4e8; color: #17603c; }}
  .badge.review {{ background: #fff0d5; color: #8a4b00; }}
  .columns {{ display: grid; grid-template-columns: repeat(auto-fit,minmax(300px,1fr)); gap: 16px; }}
  .source-page {{ width: 100%; max-height: 620px; object-fit: contain; background: #f7f8fb; border-radius: 8px; }}
  pre {{ white-space: pre-wrap; overflow-wrap: anywhere; background: #f7f8fb; padding: 14px; border-radius: 8px; max-height: 420px; overflow: auto; }}
  .warning-text,.warning {{ color: #8a4b00; }} .summary .good {{ color: #17603c; }} .muted {{ color: #526079; }}
  @media (max-width: 800px) {{ .columns {{ grid-template-columns: 1fr; }} }}
</style>
</head>
<body>
<header><h1>Book conversion review</h1><p>{html.escape(str(output_path))}</p></header>
<main>
  <div class="summary">
    <div><strong>{len(page_outputs)}</strong> pages converted</div>
    <div><strong>{uncertain_count}</strong> pages need review</div>
    <div>Words: <strong>{verification.get('source_word_count', 0)}</strong> source / <strong>{verification.get('generated_word_count', 0)}</strong> generated</div>
    <div>Numeric digits: <strong>{'preserved' if verification.get('numbers_match') else 'differences found'}</strong></div>
    {pdf_diagnostic}
    <div class="{compile_class}">{html.escape(compile_message)}</div>
    <button type="button" onclick="toggleReview(this)">Show only pages needing review</button>
  </div>
  {''.join(cards)}
</main>
<script>
let reviewOnly=false;
function toggleReview(button) {{
  reviewOnly=!reviewOnly;
  document.querySelectorAll('.page-card').forEach(card => {{
    card.style.display=(reviewOnly && card.dataset.review!=='yes')?'none':'block';
  }});
  button.textContent=reviewOnly?'Show all pages':'Show only pages needing review';
}}
</script>
</body></html>
"""


def _write_review(
    review_dir: Path,
    output_path: Path,
    page_outputs: list[PageOutput],
    compilation: dict[str, object] | None = None,
    pdf_analysis: dict[str, object] | None = None,
) -> dict[str, str]:
    review_dir.mkdir(parents=True, exist_ok=True)
    pages_dir = review_dir / "pages"
    pages_dir.mkdir(parents=True, exist_ok=True)
    for stale in pages_dir.glob("page_*_review.md"):
        stale.unlink()

    review_csv = review_dir / "review_pages.csv"
    uncertain_csv = review_dir / "review_uncertain_pages.csv"
    extracted_source_path = review_dir / "source_extracted.txt"
    extracted_source_path.write_text(
        "\n\f\n".join(page.source_text for page in page_outputs),
        encoding="utf-8",
    )
    fields = [
        "page",
        "match_ratio",
        "numbers_match",
        "missing_numbers",
        "extra_numbers",
        "uncertain",
        "safe_fallback_used",
        "warnings",
    ]

    with review_csv.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.DictWriter(fh, fieldnames=fields)
        writer.writeheader()
        for page in page_outputs:
            writer.writerow(
                {
                    "page": page.page_no,
                    "match_ratio": f"{page.match_ratio:.4f}",
                    "numbers_match": "true" if page.number_match else "false",
                    "missing_numbers": page.missing_numbers,
                    "extra_numbers": page.extra_numbers,
                    "uncertain": "true" if page.uncertain else "false",
                    "safe_fallback_used": "true" if page.used_fallback else "false",
                    "warnings": page.warnings or "",
                }
            )

    with uncertain_csv.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.DictWriter(fh, fieldnames=fields)
        writer.writeheader()
        for page in page_outputs:
            if page.uncertain:
                writer.writerow(
                    {
                        "page": page.page_no,
                        "match_ratio": f"{page.match_ratio:.4f}",
                        "numbers_match": "true" if page.number_match else "false",
                        "missing_numbers": page.missing_numbers,
                        "extra_numbers": page.extra_numbers,
                        "uncertain": "true",
                        "safe_fallback_used": "true" if page.used_fallback else "false",
                        "warnings": page.warnings or "",
                    }
                )

    uncertain_lines: list[str] = []
    for page in page_outputs:
        if page.uncertain:
            uncertain_lines.append(f"- Page {page.page_no}: text match {page.match_ratio:.2%}")
        review_md = (
            f"## Page {page.page_no}\n\n"
            "### Original text\n```text\n"
            f"{page.source_text}\n```\n\n"
            "### Generated LaTeX\n```tex\n"
            f"{page.latex_body}\n```\n\n"
            f"Text match: {page.match_ratio:.2%}\n\n"
            f"Numbers match: {'yes' if page.number_match else 'no'}\n\n"
            f"Missing numbers: {page.missing_numbers or 'none'}\n\n"
            f"Extra numbers: {page.extra_numbers or 'none'}\n\n"
            f"Needs review: {'yes' if page.uncertain else 'no'}\n\n"
            f"Warnings: {page.warnings or 'none'}\n"
        )
        (pages_dir / f"page_{page.page_no:03d}_review.md").write_text(review_md, encoding="utf-8")

    summary = [
        "# Review summary",
        "",
        f"- Pages converted: {len(page_outputs)}",
        f"- Pages needing review: {sum(1 for page in page_outputs if page.uncertain)}",
    ]
    if compilation:
        summary.append(f"- PDF created: {'yes' if compilation.get('success') else 'no'}")
        if compilation.get("message"):
            summary.append(f"- PDF note: {compilation['message']}")
        compiled_verification = compilation.get("compiled_pdf_verification") or {}
        if compiled_verification:
            summary.append(
                f"- Compiled PDF numeric verification: {compiled_verification.get('numbers_match')}"
            )
            summary.append(
                f"- Compiled PDF text similarity: {float(compiled_verification.get('text_similarity', 0.0)):.2%}"
            )
    if uncertain_lines:
        summary.extend(["", "## Pages needing review", ""])
        summary.extend(uncertain_lines)
    verification = _verification_summary(page_outputs)
    summary.extend(
        [
            "",
            "## Whole-document verification",
            "",
            f"- Original word tokens: {verification['source_word_count']}",
            f"- Generated word tokens: {verification['generated_word_count']}",
            f"- Word-count ratio: {float(verification['word_count_ratio']):.4f}",
            f"- Text similarity: {float(verification['text_similarity']):.2%}",
            f"- Original numeric tokens: {verification['source_number_count']}",
            f"- Generated numeric tokens: {verification['generated_number_count']}",
            f"- Numeric token grouping exact match: {verification['number_tokens_match_exactly']}",
            f"- Numeric digit sequence match: {verification['numeric_digits_match']}",
        ]
    )
    summary_path = review_dir / "review_summary.md"
    summary_path.write_text("\n".join(summary) + "\n", encoding="utf-8")

    verification_path = review_dir / "verification.json"
    verification_path.write_text(json.dumps(verification, indent=2), encoding="utf-8")
    verification_summary_path = review_dir / "verification_summary.md"
    verification_summary_path.write_text(
        "\n".join(
            [
                "# Verification summary",
                "",
                f"- Original word tokens: {verification['source_word_count']}",
                f"- Generated word tokens: {verification['generated_word_count']}",
                f"- Word-count ratio: {float(verification['word_count_ratio']):.4f}",
                f"- Text similarity: {float(verification['text_similarity']):.2%}",
                f"- Numeric token grouping exact match: {verification['number_tokens_match_exactly']}",
                f"- Numeric digit sequence match: {verification['numeric_digits_match']}",
                f"- Missing numbers: {', '.join(verification['missing_numbers']) or 'none'}",
                f"- Extra numbers: {', '.join(verification['extra_numbers']) or 'none'}",
                f"- Most common missing words: {', '.join(verification['missing_word_tokens']) or 'none'}",
                f"- Most common extra words: {', '.join(verification['extra_word_tokens']) or 'none'}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    report_path = review_dir / "review_report.html"
    report_path.write_text(
        _build_review_html(review_dir, output_path, page_outputs, compilation, verification, pdf_analysis),
        encoding="utf-8",
    )

    analysis_path: Path | None = None
    analysis_summary_path: Path | None = None
    if pdf_analysis:
        analysis_path = review_dir / "source_pdf_analysis.json"
        analysis_path.write_text(json.dumps(pdf_analysis, indent=2), encoding="utf-8")
        analysis_summary_path = review_dir / "source_pdf_analysis.md"
        analysis_summary_path.write_text(
            "\n".join(
                [
                    "# Source PDF analysis",
                    "",
                    f"- Pages: {pdf_analysis.get('page_count', 0)}",
                    f"- Producer: {pdf_analysis.get('producer') or 'unknown'}",
                    f"- Creator: {pdf_analysis.get('creator') or 'unknown'}",
                    f"- Creation date: {pdf_analysis.get('creation_date') or 'unknown'}",
                    f"- Text characters: {pdf_analysis.get('text_characters', 0)}",
                    f"- Embedded images: {pdf_analysis.get('embedded_images', 0)}",
                    f"- Full-page raster pages: {pdf_analysis.get('full_page_raster_pages', 0)}",
                    f"- Vector drawing objects: {pdf_analysis.get('vector_drawing_count', 0)}",
                    f"- Likely full-page raster document: {pdf_analysis.get('likely_full_page_raster_document', False)}",
                    f"- Likely raster plus text overlay: {pdf_analysis.get('likely_raster_with_text_overlay', False)}",
                    "",
                    "## Warnings",
                    "",
                    *(
                        [f"- {warning}" for warning in pdf_analysis.get("warnings", [])]
                        or ["- none"]
                    ),
                ]
            )
            + "\n",
            encoding="utf-8",
        )

    manifest = {
        "created_at": datetime.now(timezone.utc).isoformat(),
        "output_file": str(output_path),
        "review_csv": str(review_csv),
        "uncertain_csv": str(uncertain_csv),
        "extracted_source": str(extracted_source_path),
        "pages_dir": str(pages_dir),
        "summary": str(summary_path),
        "verification": str(verification_path),
        "verification_summary": str(verification_summary_path),
        "report": str(report_path),
        "compilation": compilation,
        "pdf_analysis": str(analysis_path) if analysis_path else None,
        "pdf_analysis_summary": str(analysis_summary_path) if analysis_summary_path else None,
    }
    manifest_path = review_dir / "review_manifest.json"
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")

    return {
        "review_dir": str(review_dir),
        "review_csv": str(review_csv),
        "uncertain_csv": str(uncertain_csv),
        "extracted_source": str(extracted_source_path),
        "summary": str(summary_path),
        "verification": str(verification_path),
        "verification_summary": str(verification_summary_path),
        "report": str(report_path),
        "manifest": str(manifest_path),
        "pdf_analysis": str(analysis_path) if analysis_path else "",
        "pdf_analysis_summary": str(analysis_summary_path) if analysis_summary_path else "",
    }


def _ensure_latex_packages(
    output_path: Path,
    log_callback: LogCallback | None,
) -> str | None:
    source = output_path.read_text(encoding="utf-8", errors="replace")
    required = {
        "microtype.sty": "microtype",
        "mathtools.sty": "mathtools",
        "booktabs.sty": "booktabs",
        "enumitem.sty": "enumitem",
        "listings.sty": "listings",
        "hyperref.sty": "hyperref",
        "geometry.sty": "geometry",
    }
    if "\\usepackage{fontspec}" in source:
        required["fontspec.sty"] = "fontspec"
    if "\\usepackage{polyglossia}" in source:
        required["polyglossia.sty"] = "polyglossia"
        required["bidi.sty"] = "bidi"
        required["Amiri-Regular.ttf"] = "amiri"
    if "\\usepackage{fontspec}" not in source:
        required["lmodern.sty"] = "lm"
    if "{ctex}" in source:
        required["ctex.sty"] = "ctex"
        required["xeCJK.sty"] = "xecjk"
        required["zhnumber.sty"] = "zhnumber"
        required["FandolSong-Regular.otf"] = "fandol"
    if "\\usepackage{graphicx}" in source:
        required["graphicx.sty"] = "graphics"
        required["float.sty"] = "float"
    if "xcolor}" in source:
        required["xcolor.sty"] = "xcolor"
    if "\\usepackage{pdfpages}" in source:
        required["pdfpages.sty"] = "pdfpages"
        required["eso-pic.sty"] = "eso-pic"
        required["pdflscape.sty"] = "pdflscape"
    if "\\usepackage{pgfplots}" in source:
        required["pgfplots.sty"] = "pgfplots"
        required["tikz.sty"] = "pgf"
    optional_packages = {
        "\\usepackage{multirow}": ("multirow.sty", "multirow"),
        "\\usepackage{longtable}": ("longtable.sty", "tools"),
        "\\usepackage{adjustbox}": ("adjustbox.sty", "adjustbox"),
        "\\usepackage{multicol}": ("multicol.sty", "tools"),
        "\\usepackage{pdflscape}": ("pdflscape.sty", "pdflscape"),
        "\\usepackage{siunitx}": ("siunitx.sty", "siunitx"),
        "\\usepackage{titlesec}": ("titlesec.sty", "titlesec"),
    }
    for package_line, (style_file, package_name) in optional_packages.items():
        if package_line in source:
            required[style_file] = package_name

    kpsewhich = shutil.which("kpsewhich")
    compiler = shutil.which(
        "xelatex" if "\\usepackage{fontspec}" in source else "pdflatex"
    )
    if not kpsewhich and compiler:
        sibling = Path(compiler).with_name("kpsewhich.exe")
        if sibling.is_file():
            kpsewhich = str(sibling)
    if not kpsewhich:
        return None

    missing: list[tuple[str, str]] = []
    for style_file, package_name in required.items():
        try:
            found = subprocess.run(
                [kpsewhich, style_file],
                capture_output=True,
                timeout=15,
                check=False,
            )
        except subprocess.TimeoutExpired:
            return "The LaTeX package database did not respond. Close other LaTeX jobs and try again."
        if found.returncode != 0 or not found.stdout.strip():
            missing.append((style_file, package_name))
    if not missing:
        return None

    mpm = shutil.which("mpm")
    if not mpm and compiler:
        sibling = Path(compiler).with_name("mpm.exe")
        if sibling.is_file():
            mpm = str(sibling)
    if not mpm:
        names = ", ".join(style for style, _package in missing)
        return f"Required LaTeX packages are missing: {names}"

    for style_file, package_name in missing:
        _log(log_callback, f"Installing missing MiKTeX package: {package_name}")
        try:
            installed = subprocess.run(
                [mpm, f"--install={package_name}"],
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=180,
                check=False,
            )
        except subprocess.TimeoutExpired:
            return f"Timed out while installing the MiKTeX package {package_name}."
        if installed.returncode != 0:
            detail = (installed.stderr or installed.stdout).strip()
            return f"Could not install {style_file} ({package_name}): {detail}"
    return None


def _repair_missing_graphics(output_path: Path) -> list[str]:
    """Replace unavailable AI-invented graphics with a compilable notice."""
    source = output_path.read_text(encoding="utf-8", errors="replace")
    pattern = re.compile(
        r"\\includegraphics\s*(?:\[[^\]]*\])?\s*"
        r"\{(?P<reference>\\detokenize\{[^{}]+\}|[^{}]+)\}"
    )
    missing: list[str] = []

    def replace(match: re.Match[str]) -> str:
        raw_reference = match.group("reference").strip()
        detokenized = re.fullmatch(r"\\detokenize\{([^{}]+)\}", raw_reference)
        reference = detokenized.group(1) if detokenized else raw_reference
        if any(token in reference for token in ("\\", "#", "$")):
            return match.group(0)
        candidate = Path(reference.replace("/", os.sep))
        if not candidate.is_absolute():
            candidate = output_path.parent / candidate
        candidates = [candidate]
        if not candidate.suffix:
            candidates.extend(candidate.with_suffix(ext) for ext in (".pdf", ".png", ".jpg", ".jpeg", ".eps"))
        if any(path.is_file() for path in candidates):
            return match.group(0)
        missing.append(reference)
        return r"\fbox{\textit{Image omitted: source asset unavailable}}"

    repaired = pattern.sub(replace, source)
    if missing and repaired != source:
        output_path.write_text(repaired, encoding="utf-8")
    return list(dict.fromkeys(missing))


def _repair_tabular_columns(output_path: Path) -> list[str]:
    """Expand too-small AI-generated tabular declarations to fit their rows."""
    source = output_path.read_text(encoding="utf-8", errors="replace")
    pattern = re.compile(
        r"\\begin\{tabular\}(?P<position>\[[^\]]*\])?"
        r"\{(?P<spec>(?:[^{}]|\{[^{}]*\})*)\}"
        r"(?P<body>.*?)\\end\{tabular\}",
        re.DOTALL,
    )
    repairs: list[str] = []

    def replace(match: re.Match[str]) -> str:
        spec = match.group("spec")
        repeated_columns = sum(
            int(count) * sum(1 for char in repeated_spec if char in "lcrpmbX")
            for count, repeated_spec in re.findall(r"\*\{(\d+)\}\{([^{}]+)\}", spec)
        )
        simplified = re.sub(r"@\{[^{}]*\}", "", spec)
        simplified = re.sub(r"([pmb])\{[^{}]*\}", r"\1", simplified)
        simplified = re.sub(r"\*\{\d+\}\{[^{}]+\}", "", simplified)
        declared = repeated_columns + sum(1 for char in simplified if char in "lcrpmbX")
        maximum = 0
        for line in match.group("body").splitlines():
            if r"\\" not in line:
                continue
            maximum = max(maximum, len(re.findall(r"(?<!\\)&", line)) + 1)
        if maximum <= max(declared, 1):
            return match.group(0)
        line_number = source.count("\n", 0, match.start()) + 1
        repairs.append(
            f"Adjusted a table near LaTeX line {line_number} from {declared} to {maximum} columns"
        )
        return (
            "\\begin{tabular}"
            f"{match.group('position') or ''}"
            f"{{*{{{maximum}}}{{c}}}}"
            f"{match.group('body')}"
            "\\end{tabular}"
        )

    repaired = pattern.sub(replace, source)
    if repairs:
        output_path.write_text(repaired, encoding="utf-8")
    return repairs


def _ensure_generated_preamble(output_path: Path) -> list[str]:
    """Add standard packages when generated body commands require them."""
    source = output_path.read_text(encoding="utf-8", errors="replace")
    requirements = [
        (r"\multirow", r"\usepackage{multirow}", "multirow tables"),
        (r"\begin{longtable}", r"\usepackage{longtable}", "long tables"),
        (r"\begin{adjustbox}", r"\usepackage{adjustbox}", "adjusted boxes"),
        (r"\begin{multicols}", r"\usepackage{multicol}", "multiple columns"),
        (r"\begin{landscape}", r"\usepackage{pdflscape}", "landscape pages"),
        (r"\SI{", r"\usepackage{siunitx}", "scientific units"),
        (
            "phi(",
            r"\pgfmathdeclarefunction{phi}{1}{\pgfmathparse{exp(-#1*#1/2)/sqrt(2*pi)}}",
            "the standard normal density function",
        ),
    ]
    packages: list[str] = []
    repairs: list[str] = []
    for body_marker, package_line, description in requirements:
        if body_marker in source and package_line not in source:
            packages.append(package_line)
            repairs.append(f"Added LaTeX support for {description}")
    if packages and r"\begin{document}" in source:
        source = source.replace(
            r"\begin{document}",
            "\n".join(packages) + "\n\\begin{document}",
            1,
        )
        output_path.write_text(source, encoding="utf-8")
    return repairs


def _repair_tikz_compatibility(output_path: Path) -> list[str]:
    """Normalize two common vision-model TikZ mistakes."""
    source = output_path.read_text(encoding="utf-8", errors="replace")
    repaired = source.replace(">=stealth'", ">=stealth")
    paired_range = "{0/0, 1/1, ..., 8/8}"
    expanded_paired_range = "{" + ", ".join(f"{value}/{value}" for value in range(9)) + "}"
    paired_range_count = repaired.count(paired_range)
    repaired = repaired.replace(paired_range, expanded_paired_range)
    normal_cdf_pattern = r"0\.5\*\(1\+erf\(\\x/sqrt\(2\)\)\)"
    repaired, normal_cdf_count = re.subn(
        normal_cdf_pattern,
        r"1/(1+exp(-1.702*\\x))",
        repaired,
    )
    repaired, color_name_count = re.subn(r"\bmaroon(?=!)", "Maroon", repaired)
    repaired, lavender_count = re.subn(
        r"\blavender(?=[!,}\]])",
        "Lavender",
        repaired,
    )
    color_name_count += lavender_count
    repaired, undefined_tick_count = re.subn(r"\(\\tickx\s*,", "(0,", repaired)
    repaired, tick_key_count = re.subn(
        r",\s*[xy]tick\s+distance\s*=\s*[-+]?\d*\.?\d+",
        "",
        repaired,
    )
    repaired, undefined_series_count = re.subn(
        r"\\foreach\s+\\country/\\color\s+in\s+\{[^{}]*\}\s*"
        r"\{[^{}]*\(\\x,\\y\)[^{}]*\}",
        "% Removed an invalid TikZ series with undefined coordinates",
        repaired,
        flags=re.DOTALL,
    )
    repaired, node_label_count = re.subn(
        r"(?m)^(\s*\\node\b[^{}]*\{)"
        r"(\\textcolor\{[^{}]+\}\{\\rule\{[^{}]+\}\{[^{}]+\}\})"
        r"\}\s*([^;\r\n]+);$",
        r"\1\2 \3};",
        repaired,
    )
    repaired, comment_count = re.subn(r"(?m)^(\s*)\\%(?=\s*[A-Za-z])", r"\1%", repaired)
    repairs: list[str] = []
    if ">=stealth'" in source:
        repairs.append("Updated a legacy TikZ arrow style")
    if comment_count:
        repairs.append(f"Repaired {comment_count} TikZ/comment line(s)")
    if tick_key_count:
        repairs.append(f"Removed {tick_key_count} pgfplots-only option(s) from TikZ pictures")
    if undefined_tick_count:
        repairs.append(f"Repaired {undefined_tick_count} undefined TikZ tick position(s)")
    if color_name_count:
        repairs.append(f"Normalized {color_name_count} TikZ/xcolor name(s)")
    if undefined_series_count:
        repairs.append(f"Removed {undefined_series_count} TikZ series with undefined coordinates")
    if node_label_count:
        repairs.append(f"Repaired {node_label_count} malformed TikZ legend label(s)")
    if paired_range_count:
        repairs.append(f"Expanded {paired_range_count} invalid paired TikZ range(s)")
    if normal_cdf_count:
        repairs.append(f"Replaced {normal_cdf_count} unsupported PGF normal-CDF expression(s)")
    if repairs:
        output_path.write_text(repaired, encoding="utf-8")
    return repairs


def _repair_inner_floats(output_path: Path) -> list[str]:
    """Make figures/tables non-floating when page content is boxed to one page."""
    source = output_path.read_text(encoding="utf-8", errors="replace")
    if r"\begin{adjustbox}" not in source:
        return []
    repaired, count = re.subn(
        r"\\begin\{(figure|table)\}(?:\[[^\]]*\])?",
        r"\\begin{\1}[H]",
        source,
    )
    if count and repaired != source:
        output_path.write_text(repaired, encoding="utf-8")
        return [f"Changed {count} boxed figure/table(s) to non-floating placement"]
    return []


def _repair_text_mode_bullets(output_path: Path) -> list[str]:
    """Wrap bullet symbols in math mode inside custom item labels."""
    source = output_path.read_text(encoding="utf-8", errors="replace")
    normalized, duplicate_count = re.subn(r"\$+\\bullet\$+", r"$\\bullet$", source)
    pattern = re.compile(r"(\\item\[[^\]]*?)(?<!\$)\\bullet(?!\$)([^\]]*\])")
    repaired, bare_count = pattern.subn(r"\1$\\bullet$\2", normalized)
    count = duplicate_count + bare_count
    if count and repaired != source:
        output_path.write_text(repaired, encoding="utf-8")
        return [f"Repaired {count} custom bullet label(s)"]
    return []


def _repair_split_math_fragments(output_path: Path) -> list[str]:
    """Join equality signs and following fraction commands into one math span."""
    source = output_path.read_text(encoding="utf-8", errors="replace")
    repaired, count = re.subn(
        r"(?m)\$\s*=\s*\$\s*(\\frac[^\r\n]+)$",
        r"$= \1$",
        source,
    )
    if count and repaired != source:
        output_path.write_text(repaired, encoding="utf-8")
        return [f"Repaired {count} split mathematical fraction(s)"]
    return []


def _repair_unicode_fragments(output_path: Path) -> list[str]:
    """Replace common Unicode math characters left in AI-generated LaTeX."""
    source = output_path.read_text(encoding="utf-8", errors="replace")
    replacements = {
        "∼": r"$\sim$",
        "−": "-",
        "≤": r"$\leq$",
        "≥": r"$\geq$",
    }
    repaired = source
    count = 0
    for character, latex in replacements.items():
        occurrences = repaired.count(character)
        repaired = repaired.replace(character, latex)
        count += occurrences
    if count and repaired != source:
        output_path.write_text(repaired, encoding="utf-8")
        return [f"Converted {count} Unicode mathematical symbol(s) to LaTeX"]
    return []


def _compile_latex(output_path: Path, log_callback: LogCallback | None) -> dict[str, object]:
    source = output_path.read_text(encoding="utf-8", errors="replace")
    compiler_name = "xelatex" if "\\usepackage{fontspec}" in source else "pdflatex"
    compiler = shutil.which(compiler_name)
    log_path = output_path.with_name(f"{output_path.stem}_compile.log")
    if not compiler:
        return {
            "requested": True,
            "success": False,
            "pdf_path": None,
            "log_path": None,
            "repairs": [],
            "message": f"The required {compiler_name} compiler was not found. Install MiKTeX (Windows) or TeX Live, then try again.",
        }

    missing_graphics = _repair_missing_graphics(output_path)
    structural_repairs = _repair_tabular_columns(output_path)
    preamble_repairs = _ensure_generated_preamble(output_path)
    tikz_repairs = _repair_tikz_compatibility(output_path)
    float_repairs = _repair_inner_floats(output_path)
    bullet_repairs = _repair_text_mode_bullets(output_path)
    math_repairs = _repair_split_math_fragments(output_path)
    unicode_repairs = _repair_unicode_fragments(output_path)
    repairs = [
        *missing_graphics,
        *structural_repairs,
        *preamble_repairs,
        *tikz_repairs,
        *float_repairs,
        *bullet_repairs,
        *math_repairs,
        *unicode_repairs,
    ]
    for reference in missing_graphics:
        _log(log_callback, f"Removed unavailable image reference before PDF creation: {reference}")
    for repair in structural_repairs:
        _log(log_callback, repair)
    for repair in preamble_repairs:
        _log(log_callback, repair)
    for repair in tikz_repairs:
        _log(log_callback, repair)
    for repair in float_repairs:
        _log(log_callback, repair)
    for repair in bullet_repairs:
        _log(log_callback, repair)
    for repair in math_repairs:
        _log(log_callback, repair)
    for repair in unicode_repairs:
        _log(log_callback, repair)

    package_error = _ensure_latex_packages(output_path, log_callback)
    if package_error:
        return {
            "requested": True,
            "success": False,
            "pdf_path": None,
            "log_path": None,
            "repairs": repairs,
            "message": package_error,
        }

    _log(log_callback, "Creating the ready-to-read PDF")
    command = [
        compiler,
        *(["--disable-installer"] if "miktex" in compiler.lower() else []),
        "-interaction=nonstopmode",
        "-halt-on-error",
        output_path.name,
    ]
    combined = ""
    return_code = 1
    try:
        for _ in range(2):
            completed = subprocess.run(
                command,
                cwd=output_path.parent,
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=240,
                check=False,
            )
            combined += completed.stdout + completed.stderr
            return_code = completed.returncode
            if return_code != 0:
                break
    except subprocess.TimeoutExpired as exc:
        combined += f"\nPDF creation timed out: {exc}\n"
        return_code = 124
    pdf_path = output_path.with_suffix(".pdf")
    success = return_code == 0 and pdf_path.is_file()
    if success:
        if log_path.is_file():
            log_path.unlink()
    else:
        log_path.write_text(combined, encoding="utf-8", errors="replace")
    for suffix in (".aux", ".out", ".toc", ".log", ".xdv"):
        generated = output_path.with_suffix(suffix)
        if generated.is_file():
            generated.unlink()
    return {
        "requested": True,
        "success": success,
        "pdf_path": str(pdf_path) if success else None,
        "log_path": None if success else str(log_path),
        "repairs": repairs,
        "compiler": compiler_name,
        "message": (
            (
                f"PDF created successfully after {len(repairs)} automatic repair(s)."
                if repairs
                else "PDF created successfully."
            )
            if success
            else f"LaTeX could not create the PDF. Open {log_path.name} for details."
        ),
    }


def _write_conversion_report(
    output_path: Path,
    *,
    input_path: Path,
    converted_pages: int,
    document_language: str,
    compilation: dict[str, object] | None,
    uncertain_pages: list[int],
    page_size: str = PAGE_SIZE_SOURCE,
    page_flow: str = PAGE_FLOW_SOURCE,
    photo_handling: str = PHOTO_KEEP,
    assets: list[str] | None = None,
) -> Path:
    """Write one concise, non-technical outcome/error report per conversion."""
    report_path = output_path.with_name(f"{output_path.stem}_conversion_report.txt")
    compilation = compilation or {}
    success = bool(compilation.get("success"))
    repairs = [str(item) for item in compilation.get("repairs", [])]
    page_size_label = {
        PAGE_SIZE_SOURCE: "Match source",
        PAGE_SIZE_A4: "A4",
        PAGE_SIZE_LETTER: "Letter",
    }.get(page_size, page_size)
    page_flow_label = {
        PAGE_FLOW_COMPACT: "Compact continuous flow",
        PAGE_FLOW_SOURCE: "One source unit per output page",
    }.get(page_flow, page_flow)
    photo_label = {
        PHOTO_KEEP: "Keep photos",
        PHOTO_DESCRIBE: "Replace with descriptions",
    }.get(photo_handling, photo_handling)
    lines = [
        "BOOK TO LATEX — CONVERSION REPORT",
        "",
        f"Source: {input_path}",
        f"LaTeX: {output_path}",
        f"Compiled PDF: {compilation.get('pdf_path') or 'Not created'}",
        f"Status: {'Completed successfully' if success else 'LaTeX created; compiled PDF needs attention'}",
        f"Content units converted: {converted_pages}",
        f"Document language: {DOCUMENT_LANGUAGES.get(document_language, document_language)}",
        f"Page size: {page_size_label}",
        f"Page usage: {page_flow_label}",
        f"Natural photographs: {photo_label}",
        f"Compiler: {compilation.get('compiler') or 'Not run'}",
        "",
        str(compilation.get("message") or "PDF compilation was not requested."),
    ]
    if repairs:
        lines.extend(
            [
                "",
                "Automatic repairs made before PDF creation:",
                *[f"- {repair}" for repair in repairs],
            ]
        )
    if assets:
        lines.extend(
            [
                "",
                "Required visual assets kept with the LaTeX project (photos, logos, screenshots, or visuals that could not be recreated exactly without guessing):",
                *[f"- {asset}" for asset in assets],
            ]
        )
    if uncertain_pages:
        lines.extend(
            [
                "",
                "Content units flagged for optional review:",
                ", ".join(str(page) for page in uncertain_pages),
            ]
        )
    if compilation.get("log_path"):
        lines.extend(["", f"Technical error log: {compilation['log_path']}"])
    report_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return report_path


def _safe_file_component(value: str) -> str:
    safe = re.sub(r"[^A-Za-z0-9._-]+", "_", value).strip("._")
    return safe or "converted_book"


def _remove_generated_directory(path: Path, output_path: Path) -> None:
    """Remove only a precisely named generated directory beside the output."""
    resolved_parent = output_path.parent.resolve()
    resolved_path = path.resolve()
    if resolved_path.parent != resolved_parent:
        raise ValueError(f"Refusing to clean generated files outside the output folder: {path}")
    if resolved_path.is_dir():
        shutil.rmtree(resolved_path)


def convert_book_to_latex(
    *,
    input_path: Path,
    output_path: Path,
    provider: str = "openai",
    model: str = "",
    endpoint: str = "",
    api_key: str = "",
    start_page: int = 1,
    end_page: int = 0,
    lines_per_page: int = 70,
    max_tokens: int = 4096,
    temperature: float = 0.0,
    timeout: float = 120.0,
    retries: int = 3,
    backoff: float = 1.5,
    keep_page_files: bool = False,
    no_wrapper: bool = False,
    no_llm: bool = False,
    strict_mode: bool = False,
    match_mode: str = MATCH_MODE_EXACT,
    match_error_percent: float = 5.0,
    review_dir: Path | None = None,
    no_review: bool = False,
    use_ocr: bool = False,
    ocr_force: bool = False,
    ocr_dpi: int = 220,
    ocr_lang: str = "eng",
    document_language: str = "",
    page_size: str = PAGE_SIZE_SOURCE,
    page_flow: str = PAGE_FLOW_SOURCE,
    photo_handling: str = PHOTO_KEEP,
    preserve_graphs: bool = False,
    preserve_layout: bool = False,
    preserve_color: bool = False,
    image_only: bool = False,
    vision_mode: bool = False,
    redraw_graphs: bool = False,
    style_guide: str = "",
    compile_pdf: bool = False,
    progress_callback: ProgressCallback | None = None,
    log_callback: LogCallback | None = None,
    cancel_callback: CancelCallback | None = None,
) -> dict[str, object]:
    input_path = input_path.expanduser().resolve()
    output_path = output_path.expanduser().resolve()
    provider = provider.strip().lower()
    document_language = (document_language or ocr_lang or "eng").strip().lower()
    if provider not in {"openai", "ollama"}:
        raise ValueError("AI service must be OpenAI-compatible or Ollama")
    if not api_key and provider == "openai":
        api_key = os.environ.get("OPENAI_API_KEY", "")
    resolved_endpoint = endpoint.strip() or (
        DEFAULT_OLLAMA_ENDPOINT if provider == "ollama" else DEFAULT_OPENAI_COMPAT_ENDPOINT
    )

    if not input_path.is_file():
        raise FileNotFoundError(f"Input file not found: {input_path}")
    if output_path.suffix.lower() != ".tex":
        raise ValueError("The output filename must end in .tex")
    if input_path == output_path:
        raise ValueError("Input and output must be different files")
    if start_page < 1:
        raise ValueError("Start page must be 1 or higher")
    if lines_per_page < 1:
        raise ValueError("Lines per page must be 1 or higher")
    if max_tokens < 1 or timeout <= 0 or retries < 0 or backoff < 0:
        raise ValueError("AI token, timeout, retry, and backoff settings must be positive")
    if not 0 <= temperature <= 2:
        raise ValueError("AI creativity must be between 0 and 2")
    if not 0 <= match_error_percent <= 100:
        raise ValueError("Allowed text difference must be between 0 and 100 percent")
    if ocr_dpi < 72:
        raise ValueError("OCR quality must be at least 72 DPI")
    if match_mode not in {MATCH_MODE_EXACT, MATCH_MODE_PERCENT}:
        raise ValueError("Invalid quality-check mode")
    if document_language not in DOCUMENT_LANGUAGES:
        supported = ", ".join(DOCUMENT_LANGUAGES.values())
        raise ValueError(f"Unsupported document language. Available languages: {supported}")
    if page_size not in {PAGE_SIZE_SOURCE, PAGE_SIZE_A4, PAGE_SIZE_LETTER}:
        raise ValueError("Page size must be source, a4, or letter")
    if page_flow not in {PAGE_FLOW_COMPACT, PAGE_FLOW_SOURCE}:
        raise ValueError("Page usage must be compact or source_pages")
    if photo_handling not in {PHOTO_KEEP, PHOTO_DESCRIBE}:
        raise ValueError("Photo handling must be keep or describe")

    file_ext = input_path.suffix.lower()
    source_is_pdf = file_ext == ".pdf"
    source_is_image = file_ext in IMAGE_EXTENSIONS
    source_is_visual = source_is_pdf or source_is_image
    if image_only and not source_is_visual:
        raise ValueError("Exact page appearance is available only for PDF and image files")
    exact_visual_mode = source_is_visual and image_only
    if exact_visual_mode and page_flow != PAGE_FLOW_SOURCE:
        _log(
            log_callback,
            "Exact visual copy always keeps one source page per output page; compact flow is available in editable modes",
        )
        page_flow = PAGE_FLOW_SOURCE
    if vision_mode and not source_is_visual:
        raise ValueError("Vision-assisted conversion is available only for PDF and image files")
    if vision_mode and no_llm:
        _log(log_callback, "Vision assistance is ignored because AI conversion is disabled")
        vision_mode = False
    if redraw_graphs and not vision_mode:
        _log(log_callback, "Graph redrawing is ignored because vision assistance is disabled")
        redraw_graphs = False
    render_pages = source_is_visual and (preserve_graphs or vision_mode)
    if not no_llm and not model.strip() and not exact_visual_mode:
        raise ValueError("Enter an AI model name, or choose conversion without AI")
    if preserve_color and not render_pages:
        _log(log_callback, "Color preservation is ignored because page images are not enabled")

    def check_cancelled() -> None:
        if cancel_callback and cancel_callback():
            raise ConversionCancelled("Conversion cancelled")

    _log(log_callback, f"Starting conversion: {input_path.name}")
    check_cancelled()
    pdf_analysis: dict[str, object] | None = None
    if source_is_pdf:
        pdf_analysis = analyze_pdf(input_path)
        pages, extraction_engine = _read_pdf_pages(input_path)
        source_format = f"PDF document ({extraction_engine})"
        _log(
            log_callback,
            "PDF analysis: "
            f"{pdf_analysis['page_count']} pages, "
            f"{pdf_analysis['text_characters']} text characters, "
            f"{pdf_analysis['full_page_raster_pages']} full-page raster pages, "
            f"{pdf_analysis['vector_drawing_count']} vector objects",
        )
        for warning in pdf_analysis.get("warnings", []):
            _log(log_callback, f"PDF source warning: {warning}")
    elif source_is_image:
        pages = [""] * _image_frame_count(input_path)
        source_format = "image OCR"
    else:
        if use_ocr:
            _log(log_callback, "OCR is ignored because this file has extractable document or text content")
        pages, source_format = _read_document_pages(input_path, lines_per_page)
    _log(log_callback, f"Detected input: {source_format}")
    total = len(pages)
    if total == 0:
        raise ValueError("No pages or text were found in the selected file")
    if end_page <= 0 or end_page > total:
        end_page = total
    if start_page > end_page:
        raise ValueError(f"Start page cannot be after page {end_page}")
    selected_indices = list(range(start_page - 1, end_page))

    if document_language == "auto" or ocr_lang == "auto":
        detected_language = _detect_language_from_text(
            "\n".join(pages[index] for index in selected_indices[:12])
        )
        if not detected_language and source_is_visual and not exact_visual_mode:
            try:
                if source_is_pdf:
                    _ensure_fitz()
                    detection_doc = fitz.open(input_path)  # type: ignore[union-attr]
                    try:
                        detection_page = detection_doc.load_page(selected_indices[0])
                        detection_pix = detection_page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                        assert _PILImage is not None
                        detection_image = _PILImage.frombytes(
                            "RGB",
                            (detection_pix.width, detection_pix.height),
                            detection_pix.samples,
                        )
                    finally:
                        detection_doc.close()
                else:
                    detection_image = _load_image_frame(input_path, selected_indices[0])
                detected_language = _detect_language_from_image(detection_image)
            except Exception:  # noqa: BLE001
                detected_language = None
        detected_language = detected_language or "eng"
        document_language = detected_language
        ocr_lang = detected_language
        _log(
            log_callback,
            f"Detected document language: {DOCUMENT_LANGUAGES.get(detected_language, detected_language)}",
        )

    if (use_ocr or source_is_image) and not exact_visual_mode:
        ensure_ocr_language(ocr_lang, log_callback)

    extract_visual_assets = vision_mode and not exact_visual_mode
    assets_dir_name = f"{_safe_file_component(output_path.stem)}_assets"
    assets_dir = output_path.parent / assets_dir_name
    page_assets: list[list[tuple[Path, str]]] = [[] for _ in range(total)]
    all_assets: list[tuple[Path, str]] = []
    if extract_visual_assets:
        _remove_generated_directory(assets_dir, output_path)
        assets_dir.mkdir(parents=True, exist_ok=True)
    else:
        _remove_generated_directory(assets_dir, output_path)

    images: list[Path | None] = [None] * total
    asset_dir_name = f"{_safe_file_component(output_path.stem)}_images"
    persistent_images_dir = output_path.parent / asset_dir_name
    temporary_images: tempfile.TemporaryDirectory[str] | None = None
    persistent_page_images = render_pages and (preserve_graphs or exact_visual_mode)
    if not persistent_page_images:
        _remove_generated_directory(persistent_images_dir, output_path)
    if render_pages and not persistent_page_images:
        temporary_images = tempfile.TemporaryDirectory(prefix="book_to_latex_pages_")
        images_dir = Path(temporary_images.name)
    else:
        images_dir = persistent_images_dir
    if source_is_pdf and (use_ocr or render_pages or extract_visual_assets):
        _ensure_fitz()
        if use_ocr:
            _ensure_ocr_dependencies()
        if render_pages:
            images_dir.mkdir(parents=True, exist_ok=True)
            for stale in images_dir.glob("page_*.png"):
                stale.unlink()
        pdf_doc = fitz.open(input_path)  # type: ignore[union-attr]
        try:
            if pdf_doc.page_count != total:
                total = pdf_doc.page_count
                pages = (pages + [""] * total)[:total]
                if end_page > total:
                    end_page = total
                if start_page > end_page:
                    raise ValueError(f"Start page cannot be after page {end_page}")
                selected_indices = list(range(start_page - 1, end_page))
                images = [None] * total
            for position, idx in enumerate(selected_indices, start=1):
                check_cancelled()
                source_text = pages[idx]
                if use_ocr and (ocr_force or not source_text.strip()):
                    _progress(
                        progress_callback,
                        position,
                        len(selected_indices),
                        f"Reading scanned page {idx + 1}",
                    )
                    pages[idx] = _ocr_pdf_page(pdf_doc, idx, ocr_dpi, ocr_lang, color=True)
                if render_pages:
                    _progress(
                        progress_callback,
                        position,
                        len(selected_indices),
                        f"Saving page picture {idx + 1}",
                    )
                    images[idx] = _render_pdf_page_image(
                        pdf_doc,
                        idx,
                        dpi=ocr_dpi,
                        color=preserve_color or exact_visual_mode,
                        output_dir=images_dir,
                    )
                if extract_visual_assets:
                    extracted_assets = _extract_pdf_page_assets(
                        pdf_doc,
                        idx,
                        assets_dir,
                        assets_dir_name,
                    )
                    page_assets[idx] = extracted_assets
                    all_assets.extend(extracted_assets)
        finally:
            pdf_doc.close()

    if source_is_image:
        needs_image_ocr = not exact_visual_mode
        if needs_image_ocr:
            _ensure_ocr_dependencies()
        if render_pages:
            images_dir.mkdir(parents=True, exist_ok=True)
            for stale in images_dir.glob("page_*.png"):
                stale.unlink()
        for position, idx in enumerate(selected_indices, start=1):
            check_cancelled()
            if needs_image_ocr:
                _progress(
                    progress_callback,
                    position,
                    len(selected_indices),
                    f"Reading text from image {idx + 1}",
                )
                pages[idx] = _ocr_image_frame(input_path, idx, ocr_lang)
            if render_pages:
                _progress(
                    progress_callback,
                    position,
                    len(selected_indices),
                    f"Saving page picture {idx + 1}",
                )
                images[idx] = _render_image_frame(
                    input_path,
                    idx,
                    idx + 1,
                    color=preserve_color or exact_visual_mode,
                    output_dir=images_dir,
                )
            if extract_visual_assets:
                asset = _save_image_frame_asset(
                    input_path,
                    idx,
                    idx + 1,
                    assets_dir,
                    assets_dir_name,
                )
                page_assets[idx] = [asset]
                all_assets.append(asset)

    if exact_visual_mode:
        _log(
            log_callback,
            "Using exact page appearance; original source pages will be included directly rather than converted to editable text",
        )

    processed: list[PageOutput] = []
    for position, page_zero_index in enumerate(selected_indices, start=1):
        check_cancelled()
        page_no = page_zero_index + 1
        _progress(
            progress_callback,
            position,
            len(selected_indices),
            f"Converting page {page_no} of {end_page}",
        )
        image_path = images[page_zero_index]
        image_reference = (
            f"{asset_dir_name}/{image_path.name}" if image_path is not None else None
        )
        exact_source_reference = None
        if exact_visual_mode:
            exact_source_reference = os.path.relpath(input_path, output_path.parent).replace(
                "\\", "/"
            )
        result = _page_to_latex(
            page_no=page_zero_index,
            page_text=pages[page_zero_index],
            source_format=source_format,
            style_guide=style_guide,
            vision_mode=vision_mode,
            redraw_graphs=redraw_graphs,
            match_mode=match_mode,
            match_error_percent=match_error_percent,
            model=model.strip(),
            endpoint=resolved_endpoint,
            api_key=api_key,
            max_tokens=max_tokens,
            temperature=temperature,
            timeout=timeout,
            retries=retries,
            backoff=backoff,
            strict_mode=strict_mode,
            no_llm=no_llm,
            preserve_graphs=preserve_graphs,
            preserve_layout=preserve_layout,
            image_path=image_path,
            image_reference=image_reference,
            exact_visual_mode=exact_visual_mode,
            log_callback=log_callback,
            document_language=document_language,
            exact_source_reference=exact_source_reference,
            page_size=page_size,
            page_flow=page_flow,
            preserve_color=preserve_color,
            photo_handling=photo_handling,
            page_asset_references=[reference for _path, reference in page_assets[page_zero_index]],
        )
        processed.append(result)
        _progress(
            progress_callback,
            position,
            len(selected_indices),
            f"Finished page {page_no} of {end_page}",
        )

    check_cancelled()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    if exact_visual_mode and source_is_pdf:
        body_separator = "\n"
    elif page_flow == PAGE_FLOW_COMPACT:
        body_separator = "\n\n"
    else:
        body_separator = "\n\\newpage\n"
    body = body_separator.join(page.latex_body for page in processed)
    kept_assets = _cleanup_unreferenced_assets(assets_dir, all_assets, body)
    if no_wrapper:
        output_text = body + "\n"
    else:
        header = _build_header(
            preserve_layout=preserve_layout,
            preserve_color=preserve_color,
            preserve_graphs=preserve_graphs,
            exact_visual_mode=exact_visual_mode,
            source_is_pdf=source_is_pdf,
            document_language=document_language,
            page_size=page_size,
            page_flow=page_flow,
            source_page_size_points=(
                list(pdf_analysis.get("page_size_points") or [])
                if pdf_analysis
                else None
            ),
            has_assets=bool(kept_assets),
            redraw_graphs=redraw_graphs,
            images_dir=None,
        )
        output_text = f"{header}\n{body}{_build_footer()}\n"
    temporary_output = output_path.with_suffix(".tex.tmp")
    temporary_output.write_text(output_text, encoding="utf-8")
    temporary_output.replace(output_path)

    page_files_dir: Path | None = None
    if keep_page_files:
        page_files_dir = output_path.parent / f"{_safe_file_component(output_path.stem)}_pages"
        page_files_dir.mkdir(parents=True, exist_ok=True)
        for pattern in ("page_*.tex", "page_*_source.txt"):
            for stale in page_files_dir.glob(pattern):
                stale.unlink()
        for page in processed:
            (page_files_dir / f"page_{page.page_no:03d}.tex").write_text(
                page.latex_body,
                encoding="utf-8",
                errors="replace",
            )
            (page_files_dir / f"page_{page.page_no:03d}_source.txt").write_text(
                page.source_text,
                encoding="utf-8",
                errors="replace",
            )

    compilation: dict[str, object] | None = None
    if compile_pdf:
        if no_wrapper:
            compilation = {
                "requested": True,
                "success": False,
                "pdf_path": None,
                "log_path": None,
                "message": "PDF creation is unavailable when 'LaTeX content only' is enabled.",
            }
        else:
            compilation = _compile_latex(output_path, log_callback)
            if compilation.get("success") and compilation.get("pdf_path"):
                compiled_pdf_path = Path(str(compilation["pdf_path"]))
                compilation["output_pdf_analysis"] = analyze_pdf(compiled_pdf_path)
                if not exact_visual_mode:
                    compilation["compiled_pdf_verification"] = _verify_compiled_pdf(
                        compiled_pdf_path,
                        processed,
                    )

    uncertain_pages = [page.page_no for page in processed if page.uncertain]
    review = None
    if not no_review:
        review_path = review_dir or output_path.parent / f"{_safe_file_component(output_path.stem)}_review"
        review = _write_review(review_path, output_path, processed, compilation, pdf_analysis)
    else:
        default_review_path = output_path.parent / f"{_safe_file_component(output_path.stem)}_review"
        _remove_generated_directory(default_review_path, output_path)

    if not keep_page_files:
        default_pages_path = output_path.parent / f"{_safe_file_component(output_path.stem)}_pages"
        _remove_generated_directory(default_pages_path, output_path)

    report_path = _write_conversion_report(
        output_path,
        input_path=input_path,
        converted_pages=len(processed),
        document_language=document_language,
        compilation=compilation,
        uncertain_pages=uncertain_pages,
        page_size=page_size,
        page_flow=page_flow,
        photo_handling=photo_handling,
        assets=kept_assets,
    )
    if temporary_images is not None:
        temporary_images.cleanup()

    _log(log_callback, f"Finished: pages {start_page}-{end_page} written to {output_path}")
    return {
        "input_path": str(input_path),
        "input_format": source_format,
        "pdf_analysis": pdf_analysis,
        "output_path": str(output_path),
        "pdf_path": compilation.get("pdf_path") if compilation else None,
        "images_dir": str(images_dir) if persistent_page_images else None,
        "assets_dir": str(assets_dir) if kept_assets else None,
        "assets": kept_assets,
        "page_files_dir": str(page_files_dir) if page_files_dir else None,
        "compilation": compilation,
        "start_page": start_page,
        "end_page": end_page,
        "converted_pages": len(processed),
        "uncertain_count": len(uncertain_pages),
        "uncertain_pages": uncertain_pages,
        "review": review,
        "report_path": str(report_path),
        "document_language": document_language,
        "page_size": page_size,
        "page_flow": page_flow,
        "photo_handling": photo_handling,
        "strict_mode": strict_mode,
        "preserve_graphs": preserve_graphs,
        "preserve_layout": preserve_layout,
        "preserve_color": preserve_color,
        "exact_visual_mode": exact_visual_mode,
        "vision_mode": vision_mode,
        "redraw_graphs": redraw_graphs,
    }


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Convert book pages to LaTeX")
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--provider", default="openai", choices=["openai", "ollama"])
    parser.add_argument("--model", default="")
    parser.add_argument("--endpoint", default="")
    parser.add_argument("--api-key", default="")
    parser.add_argument("--start-page", type=int, default=1)
    parser.add_argument("--end-page", type=int, default=0)
    parser.add_argument("--lines-per-page", type=int, default=70)
    parser.add_argument("--max-tokens", type=int, default=4096)
    parser.add_argument("--temperature", type=float, default=0.0)
    parser.add_argument("--timeout", type=float, default=120.0)
    parser.add_argument("--retries", type=int, default=3)
    parser.add_argument("--backoff", type=float, default=1.5)
    parser.add_argument("--keep-page-files", action="store_true")
    parser.add_argument("--no-wrapper", action="store_true")
    parser.add_argument("--no-llm", action="store_true", help="Skip LLM and do local conversion")
    parser.add_argument("--strict-mode", action="store_true")
    parser.add_argument("--match-mode", default=MATCH_MODE_EXACT, choices=[MATCH_MODE_EXACT, MATCH_MODE_PERCENT])
    parser.add_argument("--match-error-percent", type=float, default=5.0)
    parser.add_argument("--review-dir", default="")
    parser.add_argument("--no-review", action="store_true")
    parser.add_argument("--ocr", dest="use_ocr", action="store_true", help="Enable OCR for PDF")
    parser.add_argument("--ocr-force", action="store_true", help="OCR every PDF page")
    parser.add_argument("--ocr-dpi", type=int, default=220)
    parser.add_argument("--ocr-lang", default="eng")
    parser.add_argument(
        "--document-language",
        default="eng",
        choices=sorted(DOCUMENT_LANGUAGES),
        help="Document language (eng or ara); Arabic automatically uses XeLaTeX and RTL layout",
    )
    parser.add_argument(
        "--page-size",
        default=PAGE_SIZE_SOURCE,
        choices=[PAGE_SIZE_SOURCE, PAGE_SIZE_A4, PAGE_SIZE_LETTER],
        help="Use the source page size, A4, or US Letter",
    )
    parser.add_argument(
        "--page-flow",
        default=PAGE_FLOW_SOURCE,
        choices=[PAGE_FLOW_COMPACT, PAGE_FLOW_SOURCE],
        help="Flow content compactly or keep every source unit on its own output page",
    )
    parser.add_argument(
        "--photo-handling",
        default=PHOTO_KEEP,
        choices=[PHOTO_KEEP, PHOTO_DESCRIBE],
        help="Keep natural photographs or replace them with written descriptions",
    )
    parser.add_argument("--preserve-graphs", action="store_true", help="Embed rendered page images")
    parser.add_argument("--preserve-layout", action="store_true", help="Prefer page layout preservation")
    parser.add_argument("--preserve-color", action="store_true", help="Keep color in rendered page images")
    parser.add_argument("--image-only", action="store_true", help="Use full-page images for exact PDF appearance")
    parser.add_argument("--vision", action="store_true", help="Send rendered PDF/image pages to a vision-capable AI model")
    parser.add_argument("--redraw-graphs", action="store_true", help="Reconstruct graphs, tables and technical visuals as semantic LaTeX")
    parser.add_argument("--style-guide", default="", help="Optional UTF-8 text/Markdown file containing project LaTeX conventions")
    parser.add_argument("--compile-pdf", action="store_true", help="Also create a ready-to-read PDF")
    return parser


def main() -> None:
    parser = _build_parser()
    args = parser.parse_args()

    review_path = Path(args.review_dir).expanduser().resolve() if args.review_dir else None
    style_guide = ""
    if args.style_guide:
        style_guide = Path(args.style_guide).expanduser().read_text(
            encoding="utf-8", errors="replace"
        )

    try:
        result = convert_book_to_latex(
            input_path=Path(args.input),
            output_path=Path(args.output),
            provider=args.provider,
            model=args.model.strip(),
            endpoint=args.endpoint.strip(),
            api_key=args.api_key.strip(),
            start_page=args.start_page,
            end_page=args.end_page,
            lines_per_page=args.lines_per_page,
            max_tokens=args.max_tokens,
            temperature=args.temperature,
            timeout=args.timeout,
            retries=args.retries,
            backoff=args.backoff,
            keep_page_files=args.keep_page_files,
            no_wrapper=args.no_wrapper,
            no_llm=args.no_llm,
            strict_mode=args.strict_mode,
            match_mode=args.match_mode,
            match_error_percent=args.match_error_percent,
            review_dir=review_path,
            no_review=args.no_review,
            use_ocr=args.use_ocr,
            ocr_force=args.ocr_force,
            ocr_dpi=args.ocr_dpi,
            ocr_lang=args.ocr_lang,
            document_language=args.document_language,
            page_size=args.page_size,
            page_flow=args.page_flow,
            photo_handling=args.photo_handling,
            preserve_graphs=args.preserve_graphs,
            preserve_layout=args.preserve_layout,
            preserve_color=args.preserve_color,
            image_only=args.image_only,
            vision_mode=args.vision,
            redraw_graphs=args.redraw_graphs,
            style_guide=style_guide,
            compile_pdf=args.compile_pdf,
            progress_callback=lambda current, total, message: print(f"[{current}/{total}] {message}"),
            log_callback=lambda message: print(message),
        )
    except Exception as exc:  # noqa: BLE001
        print(f"Conversion failed: {exc}", file=sys.stderr)
        raise SystemExit(1) from exc

    print(f"Converted {result['start_page']}-{result['end_page']} -> {result['output_path']}")
    print(f"Uncertain pages: {result['uncertain_count']}")
    if result.get("pdf_path"):
        print(f"PDF: {result['pdf_path']}")
    elif result.get("compilation"):
        print(f"PDF note: {result['compilation']['message']}")
    if result.get("review"):
        print(f"Review package: {result['review']['review_dir']}")


if __name__ == "__main__":
    main()
