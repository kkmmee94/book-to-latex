from __future__ import annotations

import csv
import json
import shutil
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

import fitz
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from book_to_latex import (
    APP_VERSION,
    DEFAULT_OLLAMA_ENDPOINT,
    PAGE_FLOW_COMPACT,
    PAGE_FLOW_SOURCE,
    PAGE_SIZE_A4,
    PHOTO_DESCRIBE,
    PHOTO_KEEP,
    ConversionCancelled,
    _compile_latex,
    _detect_language_from_text,
    _escape_latex,
    _latex_to_plain,
    _numeric_comparison,
    _post_validate_latex,
    _query_llm,
    _read_document_pages,
    analyze_pdf,
    check_for_updates,
    convert_book_to_latex,
    ensure_ocr_language,
    runtime_capabilities,
)

SAMPLE_TEXT = "Title & cost 50% — 5°\n\nA_b #1 and {braces}"


class UpdateTests(unittest.TestCase):
    def test_update_check_compares_public_release_version(self) -> None:
        class FakeResponse:
            def __enter__(self):
                return self

            def __exit__(self, *_args):
                return False

            def read(self) -> bytes:
                return json.dumps(
                    {
                        "tag_name": "v99.0.0",
                        "html_url": "https://github.com/kkmmee94/book-to-latex/releases/tag/v99.0.0",
                        "assets": [
                            {
                                "name": "Book-to-LaTeX-Setup.exe",
                                "browser_download_url": "https://example.invalid/installer.exe",
                            }
                        ],
                    }
                ).encode("utf-8")

        with patch("book_to_latex.urllib.request.urlopen", return_value=FakeResponse()):
            result = check_for_updates()
        self.assertEqual(result["current_version"], APP_VERSION)
        self.assertEqual(result["latest_version"], "99.0.0")
        self.assertTrue(result["update_available"])
        self.assertEqual(result["assets"][0]["name"], "Book-to-LaTeX-Setup.exe")


class TextConversionTests(unittest.TestCase):
    def test_chinese_text_is_detected_without_an_english_assumption(self) -> None:
        self.assertEqual(_detect_language_from_text("这是一个中文讲义测试。统计学课程。"), "chi_sim")

    def test_missing_ocr_language_is_downloaded_to_user_data(self) -> None:
        class FakeLanguageResponse:
            def __enter__(self):
                return self

            def __exit__(self, *_args):
                return False

            def read(self, _size: int = -1) -> bytes:
                return b"trained-language-data" * 10_000

        with tempfile.TemporaryDirectory() as temporary, patch(
            "book_to_latex._user_tessdata_dir",
            return_value=Path(temporary),
        ), patch(
            "book_to_latex.pytesseract.get_languages",
            return_value=["eng", "osd"],
        ), patch(
            "book_to_latex.urllib.request.urlopen",
            return_value=FakeLanguageResponse(),
        ):
            ready = ensure_ocr_language("chi_sim")
            installed = Path(temporary) / "chi_sim.traineddata"
            self.assertEqual(ready, ["chi_sim"])
            self.assertTrue(installed.is_file())
            self.assertGreater(installed.stat().st_size, 100_000)

    def test_arabic_ocr_data_is_available_to_the_app(self) -> None:
        self.assertIn("ara", runtime_capabilities()["ocr_languages"])

    def test_plain_text_escaping_is_valid_and_single_escaped(self) -> None:
        escaped = _escape_latex(SAMPLE_TEXT)
        self.assertIn(r"\&", escaped)
        self.assertNotIn(r"\\&", escaped)
        self.assertIn(r"\%", escaped)
        self.assertIn(r"\_", escaped)
        self.assertIn(r"\textdegree{}", escaped)

    def test_model_output_repairs_plain_specials_but_keeps_alignment(self) -> None:
        repaired, warnings = _post_validate_latex(
            "A & B cost 50% #1\n"
            "\\begin{align*}\n"
            "x &= 1 \\\\\n"
            "\\end{align*}"
        )
        self.assertIn(r"A \& B cost 50\% \#1", repaired)
        self.assertIn("x &= 1", repaired)
        self.assertTrue(any("ampersand" in warning for warning in warnings))

    def test_latex_text_extraction_keeps_formatting_arguments(self) -> None:
        visible = _latex_to_plain(r"\textbf{Hello} \emph{world} \& 50\%")
        self.assertEqual(visible, "Hello world & 50%")

    def test_literal_backslash_survives_review_extraction(self) -> None:
        visible = _latex_to_plain(r"C:\textbackslash{}Users\textbackslash{}Books")
        self.assertEqual(visible, r"C:\Users\Books")

    def test_local_conversion_matches_and_creates_review(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "source.txt"
            output = root / "finished.tex"
            source.write_text(SAMPLE_TEXT, encoding="utf-8")
            result = convert_book_to_latex(
                input_path=source,
                output_path=output,
                no_llm=True,
            )
            self.assertEqual(result["uncertain_count"], 0)
            self.assertTrue(output.is_file())
            self.assertIn(r"Title \& cost 50\%", output.read_text(encoding="utf-8"))
            report = Path(str(result["review"]["report"]))
            self.assertTrue(report.is_file())
            self.assertIn("Book conversion review", report.read_text(encoding="utf-8"))

    def test_compact_and_source_page_flows_are_both_preserved(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "source.txt"
            source.write_text("First unit\nSecond unit", encoding="utf-8")
            compact_output = root / "compact.tex"
            source_pages_output = root / "source-pages.tex"
            compact = convert_book_to_latex(
                input_path=source,
                output_path=compact_output,
                no_llm=True,
                lines_per_page=1,
                page_flow=PAGE_FLOW_COMPACT,
                page_size=PAGE_SIZE_A4,
                no_review=True,
            )
            separated = convert_book_to_latex(
                input_path=source,
                output_path=source_pages_output,
                no_llm=True,
                lines_per_page=1,
                page_flow=PAGE_FLOW_SOURCE,
                page_size=PAGE_SIZE_A4,
                no_review=True,
            )
            self.assertNotIn(r"\newpage", compact_output.read_text(encoding="utf-8"))
            self.assertIn(r"\newpage", source_pages_output.read_text(encoding="utf-8"))
            self.assertEqual(compact["page_flow"], PAGE_FLOW_COMPACT)
            self.assertEqual(separated["page_flow"], PAGE_FLOW_SOURCE)

    def test_pdf_compilation_when_pdflatex_is_available(self) -> None:
        if not shutil.which("pdflatex"):
            self.skipTest("pdflatex is not installed")
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "source.txt"
            output = root / "finished.tex"
            source.write_text(SAMPLE_TEXT, encoding="utf-8")
            result = convert_book_to_latex(
                input_path=source,
                output_path=output,
                no_llm=True,
                compile_pdf=True,
                no_review=True,
            )
            self.assertTrue(result["compilation"]["success"])
            self.assertTrue(Path(str(result["pdf_path"])).is_file())
            self.assertTrue(
                result["compilation"]["compiled_pdf_verification"]["numbers_match"]
            )
            self.assertTrue(Path(str(result["report_path"])).is_file())
            self.assertIsNone(result["compilation"]["log_path"])

    def test_missing_ai_graphic_is_omitted_reported_and_compiles(self) -> None:
        if not shutil.which("pdflatex"):
            self.skipTest("pdflatex is not installed")
        with tempfile.TemporaryDirectory() as temporary, patch(
            "book_to_latex._query_llm",
            return_value="Before\n\\includegraphics{invented_logo.png}\nAfter",
        ):
            root = Path(temporary)
            source = root / "source.txt"
            output = root / "finished.tex"
            source.write_text("Before After", encoding="utf-8")
            result = convert_book_to_latex(
                input_path=source,
                output_path=output,
                model="test-model",
                compile_pdf=True,
                no_review=True,
            )
            self.assertTrue(result["compilation"]["success"])
            self.assertEqual(result["compilation"]["repairs"], ["invented_logo.png"])
            self.assertNotIn("invented_logo.png", output.read_text(encoding="utf-8"))
            report = Path(str(result["report_path"])).read_text(encoding="utf-8")
            self.assertIn("invented_logo.png", report)

    def test_compile_repairs_table_package_and_tikz_mistakes_in_one_pass(self) -> None:
        if not shutil.which("pdflatex"):
            self.skipTest("pdflatex is not installed")
        with tempfile.TemporaryDirectory() as temporary:
            output = Path(temporary) / "repairs.tex"
            output.write_text(
                r"""\documentclass{article}
\usepackage{tikz}
\begin{document}
\begin{tabular}{c|c}
0 & 20 & 40 \\
\multirow{2}{*}{Asia} & 1 & 2 \\
\end{tabular}
\begin{tikzpicture}[>=stealth', xtick distance=0.5, ytick distance=0.2]
\% Axes
\draw[->] (0,0) -- (1,1);
\end{tikzpicture}
\end{document}
""",
                encoding="utf-8",
            )
            result = _compile_latex(output, None)
            self.assertTrue(result["success"], result)
            self.assertTrue(any("table" in repair for repair in result["repairs"]))
            self.assertTrue(any("multirow" in repair for repair in result["repairs"]))
            self.assertTrue(any("TikZ" in repair for repair in result["repairs"]))

    def test_compile_repairs_floats_inside_one_page_boxes(self) -> None:
        if not shutil.which("pdflatex"):
            self.skipTest("pdflatex is not installed")
        with tempfile.TemporaryDirectory() as temporary:
            output = Path(temporary) / "boxed-float.tex"
            output.write_text(
                r"""\documentclass{article}
\usepackage{float}
\usepackage{adjustbox}
\begin{document}
\begin{adjustbox}{max width=\textwidth,max totalheight=.9\textheight}
\begin{minipage}{\textwidth}
\begin{figure}[h]
\centering\rule{3cm}{2cm}
\end{figure}
\end{minipage}
\end{adjustbox}
\end{document}
""",
                encoding="utf-8",
            )
            result = _compile_latex(output, None)
            self.assertTrue(result["success"], result)
            self.assertTrue(any("non-floating" in repair for repair in result["repairs"]))

    def test_arabic_uses_rtl_xelatex_and_compiles(self) -> None:
        if not shutil.which("xelatex"):
            self.skipTest("xelatex is not installed")
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "arabic.txt"
            output = root / "arabic.tex"
            source.write_text("مرحبا بالعالم", encoding="utf-8")
            result = convert_book_to_latex(
                input_path=source,
                output_path=output,
                no_llm=True,
                document_language="ara",
                compile_pdf=True,
                no_review=True,
            )
            tex = output.read_text(encoding="utf-8")
            self.assertIn(r"\setdefaultlanguage{arabic}", tex)
            self.assertEqual(result["compilation"]["compiler"], "xelatex")
            self.assertTrue(result["compilation"]["success"], result["compilation"])

    def test_chinese_uses_unicode_xelatex_when_ctex_is_available(self) -> None:
        if not shutil.which("xelatex") or not shutil.which("kpsewhich"):
            self.skipTest("XeLaTeX is not installed")
        ctex = subprocess.run(
            ["kpsewhich", "ctex.sty"],
            capture_output=True,
            text=True,
            timeout=20,
            check=False,
        )
        if ctex.returncode != 0 or not ctex.stdout.strip():
            self.skipTest("CTeX is not installed")
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "chinese.txt"
            output = root / "chinese.tex"
            source.write_text("统计学课程：这是一个中文讲义测试。", encoding="utf-8")
            result = convert_book_to_latex(
                input_path=source,
                output_path=output,
                no_llm=True,
                document_language="chi_sim",
                compile_pdf=True,
                no_review=True,
            )
            self.assertTrue(result["compilation"]["success"], result["compilation"])
            tex = output.read_text(encoding="utf-8")
            self.assertIn("{ctex}", tex)

    def test_ai_failure_uses_safe_fallback_and_marks_review(self) -> None:
        with tempfile.TemporaryDirectory() as temporary, patch(
            "book_to_latex._query_llm",
            side_effect=RuntimeError("service unavailable"),
        ):
            root = Path(temporary)
            source = root / "source.txt"
            output = root / "finished.tex"
            source.write_text("Hello world", encoding="utf-8")
            result = convert_book_to_latex(
                input_path=source,
                output_path=output,
                model="test-model",
            )
            self.assertEqual(result["uncertain_pages"], [1])
            with Path(str(result["review"]["review_csv"])).open(
                "r", encoding="utf-8", newline=""
            ) as stream:
                rows = list(csv.DictReader(stream))
            self.assertEqual(rows[0]["safe_fallback_used"], "true")
            self.assertIn("safe local conversion", rows[0]["warnings"])

    def test_changed_number_is_always_flagged(self) -> None:
        with tempfile.TemporaryDirectory() as temporary, patch(
            "book_to_latex._query_llm",
            return_value="The reported value is 43.",
        ):
            root = Path(temporary)
            source = root / "source.txt"
            output = root / "finished.tex"
            source.write_text("The reported value is 42.", encoding="utf-8")
            result = convert_book_to_latex(
                input_path=source,
                output_path=output,
                model="test-model",
                match_mode="percentage",
                match_error_percent=50,
            )
            self.assertEqual(result["uncertain_pages"], [1])
            with Path(str(result["review"]["review_csv"])).open(
                "r", encoding="utf-8", newline=""
            ) as stream:
                row = next(csv.DictReader(stream))
            self.assertEqual(row["numbers_match"], "false")
            self.assertEqual(row["missing_numbers"], "42")
            self.assertEqual(row["extra_numbers"], "43")

    def test_formula_layout_grouping_does_not_create_false_numeric_failure(self) -> None:
        comparison = _numeric_comparison("S 21 X 2", "S 2 1 X 2")
        self.assertFalse(comparison["number_tokens_match_exactly"])
        self.assertTrue(comparison["numeric_digits_match"])
        self.assertTrue(comparison["numbers_match"])

    def test_review_rerun_removes_stale_page_files(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "source.txt"
            output = root / "finished.tex"
            source.write_text("First page", encoding="utf-8")
            first = convert_book_to_latex(input_path=source, output_path=output, no_llm=True)
            pages_dir = Path(str(first["review"]["review_dir"])) / "pages"
            stale = pages_dir / "page_999_review.md"
            stale.write_text("stale", encoding="utf-8")
            convert_book_to_latex(input_path=source, output_path=output, no_llm=True)
            self.assertFalse(stale.exists())

    def test_cancellation_stops_before_output(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "source.txt"
            output = root / "finished.tex"
            source.write_text("Hello", encoding="utf-8")
            with self.assertRaises(ConversionCancelled):
                convert_book_to_latex(
                    input_path=source,
                    output_path=output,
                    no_llm=True,
                    cancel_callback=lambda: True,
                )
            self.assertFalse(output.exists())


class PdfImageTests(unittest.TestCase):
    @staticmethod
    def _make_pdf(path: Path, pages: int = 3) -> None:
        document = fitz.open()
        for page_number in range(1, pages + 1):
            page = document.new_page()
            page.insert_text((72, 72), f"Test page {page_number}")
        document.save(str(path))
        document.close()

    def test_exact_pdf_mode_reuses_source_pdf_without_page_images(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "source.pdf"
            output = root / "My finished book.tex"
            self._make_pdf(source)
            stale_images = root / "My_finished_book_images"
            stale_images.mkdir()
            (stale_images / "stale.png").write_bytes(b"stale")
            result = convert_book_to_latex(
                input_path=source,
                output_path=output,
                image_only=True,
                page_flow=PAGE_FLOW_COMPACT,
                start_page=2,
                end_page=2,
                compile_pdf=bool(shutil.which("pdflatex")),
                no_review=True,
            )
            tex = output.read_text(encoding="utf-8")
            self.assertIn(r"\includepdf[pages={2},fitpaper=true", tex)
            self.assertIn(r"\detokenize{source.pdf}", tex)
            self.assertIsNone(result["images_dir"])
            self.assertEqual(result["page_flow"], PAGE_FLOW_SOURCE)
            self.assertFalse((root / "My_finished_book_images").exists())
            if shutil.which("pdflatex"):
                self.assertTrue(result["compilation"]["success"])
                compiled_analysis = analyze_pdf(Path(str(result["pdf_path"])))
                self.assertEqual(compiled_analysis["page_count"], 1)
                source_analysis = analyze_pdf(source)
                for actual, expected in zip(
                    compiled_analysis["page_size_points"],
                    source_analysis["page_size_points"],
                    strict=True,
                ):
                    self.assertAlmostEqual(actual, expected, delta=1.0)

    def test_exact_copy_can_be_placed_on_a4_without_losing_source_page_mode(self) -> None:
        if not shutil.which("pdflatex"):
            self.skipTest("pdflatex is not installed")
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "landscape.pdf"
            document = fitz.open()
            page = document.new_page(width=400, height=250)
            page.insert_text((40, 80), "Landscape source", fontsize=24)
            document.save(source)
            document.close()
            result = convert_book_to_latex(
                input_path=source,
                output_path=root / "a4-copy.tex",
                image_only=True,
                page_size=PAGE_SIZE_A4,
                compile_pdf=True,
                no_review=True,
            )
            self.assertTrue(result["compilation"]["success"])
            tex = (root / "a4-copy.tex").read_text(encoding="utf-8")
            self.assertIn("a4paper", tex)
            self.assertIn("fitpaper=false", tex)
            analysis = analyze_pdf(Path(str(result["pdf_path"])))
            self.assertAlmostEqual(analysis["page_size_points"][0], 595.276, delta=1.0)
            self.assertAlmostEqual(analysis["page_size_points"][1], 841.89, delta=1.0)

    def test_close_layout_keeps_only_referenced_real_photo_assets(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            photo = root / "photo.png"
            photo_document = fitz.open()
            photo_page = photo_document.new_page(width=180, height=120)
            photo_page.draw_rect(photo_page.rect, color=(0.1, 0.4, 0.8), fill=(0.1, 0.4, 0.8))
            photo_page.insert_text((25, 65), "BEAR PHOTO", fontsize=18, color=(1, 1, 1))
            photo_page.get_pixmap(alpha=False).save(photo)
            photo_document.close()

            source = root / "source-with-photo.pdf"
            document = fitz.open()
            page = document.new_page(width=600, height=800)
            page.insert_text((60, 70), "Wildlife report", fontsize=20)
            page.insert_image(fitz.Rect(80, 130, 440, 370), filename=str(photo))
            document.save(source)
            document.close()

            def include_exact_asset(**kwargs):
                references = kwargs["page_asset_references"]
                self.assertTrue(references)
                return (
                    "\\section*{Wildlife report}\n"
                    "\\begin{figure}[H]\\centering\n"
                    f"\\includegraphics[width=.6\\linewidth]{{{references[0]}}}\n"
                    "\\caption{Bear in its habitat}\\end{figure}"
                )

            def classify_as_photo(**kwargs):
                return [
                    {
                        "reference": kwargs["asset_references"][0],
                        "kind": "photo",
                        "description": "A bear in its habitat",
                    }
                ]

            with patch("book_to_latex._query_llm", side_effect=include_exact_asset), patch(
                "book_to_latex._query_visual_inventory",
                side_effect=classify_as_photo,
            ):
                result = convert_book_to_latex(
                    input_path=source,
                    output_path=root / "photo-project.tex",
                    model="vision-test",
                    vision_mode=True,
                    redraw_graphs=True,
                    photo_handling=PHOTO_KEEP,
                    no_review=True,
                )
            self.assertIsNone(result["images_dir"])
            self.assertTrue(result["assets"])
            assets_dir = Path(str(result["assets_dir"]))
            self.assertTrue(assets_dir.is_dir())
            self.assertEqual(len(list(assets_dir.iterdir())), 1)
            tex = (root / "photo-project.tex").read_text(encoding="utf-8")
            self.assertIn(result["assets"][0], tex)
            self.assertIn(r"\usepackage{graphicx}", tex)

    def test_raster_graph_is_retried_as_semantic_latex_and_asset_is_removed(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            graph = root / "graph.png"
            graph_document = fitz.open()
            graph_page = graph_document.new_page(width=220, height=150)
            graph_page.draw_line((20, 125), (200, 125), width=2)
            graph_page.draw_line((20, 125), (20, 20), width=2)
            graph_page.draw_polyline([(25, 110), (70, 80), (115, 95), (170, 35)], color=(0, 0.3, 0.9), width=3)
            graph_page.get_pixmap(alpha=False).save(graph)
            graph_document.close()
            source = root / "graph-source.pdf"
            document = fitz.open()
            page = document.new_page(width=600, height=800)
            page.insert_text((70, 70), "Results over time", fontsize=20)
            page.insert_image(fitz.Rect(80, 140, 520, 440), filename=str(graph))
            document.save(source)
            document.close()

            call_count = 0

            def classify_graph(**kwargs):
                return [
                    {
                        "reference": kwargs["asset_references"][0],
                        "kind": "graph",
                        "description": "A line graph of results over time",
                        "reconstructable": "yes",
                    }
                ]

            def retry_semantically(**kwargs):
                nonlocal call_count
                call_count += 1
                if not kwargs["force_semantic_retry"]:
                    return f"\\includegraphics{{{kwargs['semantic_asset_references'][0]}}}"
                return (
                    "\\begin{tikzpicture}\n"
                    "\\begin{axis}[xlabel={Time},ylabel={Result}]\n"
                    "\\addplot coordinates {(1,1) (2,3) (3,2) (4,5)};\n"
                    "\\end{axis}\n\\end{tikzpicture}"
                )

            with patch("book_to_latex._query_visual_inventory", side_effect=classify_graph), patch(
                "book_to_latex._query_llm",
                side_effect=retry_semantically,
            ):
                result = convert_book_to_latex(
                    input_path=source,
                    output_path=root / "semantic-graph.tex",
                    model="vision-test",
                    vision_mode=True,
                    redraw_graphs=True,
                    no_review=True,
                )
            tex = (root / "semantic-graph.tex").read_text(encoding="utf-8")
            self.assertEqual(call_count, 2)
            self.assertIn(r"\begin{axis}", tex)
            self.assertIsNone(result["assets_dir"])
            self.assertFalse((root / "semantic-graph_assets").exists())

    def test_describe_photo_policy_adds_description_and_removes_photo_asset(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            photo = root / "person.png"
            image_document = fitz.open()
            image_page = image_document.new_page(width=160, height=160)
            image_page.draw_circle((80, 50), 25, color=(0.2, 0.2, 0.2), fill=(0.8, 0.6, 0.4))
            image_page.draw_rect(fitz.Rect(45, 80, 115, 150), color=(0.1, 0.3, 0.7), fill=(0.1, 0.3, 0.7))
            image_page.get_pixmap(alpha=False).save(photo)
            image_document.close()
            source = root / "person-source.pdf"
            document = fitz.open()
            page = document.new_page(width=600, height=800)
            page.insert_text((70, 70), "Author profile", fontsize=20)
            page.insert_image(fitz.Rect(120, 140, 360, 380), filename=str(photo))
            document.save(source)
            document.close()

            def classify_photo(**kwargs):
                return [
                    {
                        "reference": kwargs["asset_references"][0],
                        "kind": "photo",
                        "description": "A portrait of the author wearing a blue shirt",
                    }
                ]

            with patch("book_to_latex._query_visual_inventory", side_effect=classify_photo), patch(
                "book_to_latex._query_llm",
                return_value="\\section*{Author profile}",
            ):
                result = convert_book_to_latex(
                    input_path=source,
                    output_path=root / "described-photo.tex",
                    model="vision-test",
                    vision_mode=True,
                    redraw_graphs=True,
                    photo_handling=PHOTO_DESCRIBE,
                    no_review=True,
                )
            tex = (root / "described-photo.tex").read_text(encoding="utf-8")
            self.assertIn("Photo description:", tex)
            self.assertIn("portrait of the author", tex)
            self.assertNotIn(r"\includegraphics", tex)
            self.assertIsNone(result["assets_dir"])

    def test_vision_page_images_are_temporary(self) -> None:
        with tempfile.TemporaryDirectory() as temporary, patch(
            "book_to_latex._query_llm", return_value="Test page 1"
        ):
            root = Path(temporary)
            source = root / "source.pdf"
            output = root / "finished.tex"
            self._make_pdf(source, pages=1)
            result = convert_book_to_latex(
                input_path=source,
                output_path=output,
                model="vision-test",
                vision_mode=True,
                no_review=True,
            )
            self.assertIsNone(result["images_dir"])
            self.assertFalse((root / "finished_images").exists())

    def test_editable_source_page_mode_cannot_spill_one_unit_to_extra_pages(self) -> None:
        if not shutil.which("pdflatex"):
            self.skipTest("pdflatex is not installed")
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            source = root / "source.pdf"
            output = root / "source-page.tex"
            self._make_pdf(source, pages=1)
            result = convert_book_to_latex(
                input_path=source,
                output_path=output,
                no_llm=True,
                preserve_layout=True,
                page_flow=PAGE_FLOW_SOURCE,
                compile_pdf=True,
                no_review=True,
            )
            self.assertTrue(result["compilation"]["success"])
            analysis = analyze_pdf(Path(str(result["pdf_path"])))
            self.assertEqual(analysis["page_count"], 1)
            source_analysis = analyze_pdf(source)
            for actual, expected in zip(
                analysis["page_size_points"],
                source_analysis["page_size_points"],
                strict=True,
            ):
                self.assertAlmostEqual(actual, expected, delta=0.1)
            tex = output.read_text(encoding="utf-8")
            self.assertIn(r"\begin{adjustbox}", tex)

    def test_image_input_runs_ocr_automatically(self) -> None:
        if not runtime_capabilities().get("tesseract"):
            self.skipTest("Tesseract is not installed")
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            image_path = root / "scanned-page.png"
            output = root / "finished.tex"
            document = fitz.open()
            page = document.new_page(width=900, height=300)
            page.insert_text((60, 170), "HELLO OCR 123", fontsize=52)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            pixmap.save(str(image_path))
            document.close()
            result = convert_book_to_latex(
                input_path=image_path,
                output_path=output,
                no_llm=True,
                no_review=True,
            )
            converted = output.read_text(encoding="utf-8")
            self.assertIn("HELLO OCR 123", converted)
            self.assertEqual(result["input_format"], "image OCR")

    def test_pdf_analysis_detects_full_page_raster(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            image_path = root / "page.png"
            pdf_path = root / "raster.pdf"
            image_document = fitz.open()
            image_page = image_document.new_page(width=600, height=800)
            image_page.insert_text((80, 200), "Raster source", fontsize=40)
            image_page.get_pixmap().save(str(image_path))
            image_document.close()
            document = fitz.open()
            page = document.new_page(width=600, height=800)
            page.insert_image(page.rect, filename=str(image_path))
            document.save(str(pdf_path))
            document.close()
            analysis = analyze_pdf(pdf_path)
            self.assertEqual(analysis["page_count"], 1)
            self.assertEqual(analysis["full_page_raster_pages"], 1)
            self.assertTrue(analysis["likely_full_page_raster_document"])


class OllamaIntegrationTests(unittest.TestCase):
    def test_native_ollama_payload_disables_thinking_and_carries_vision(self) -> None:
        class FakeResponse:
            def __enter__(self):
                return self

            def __exit__(self, *_args):
                return False

            def read(self) -> bytes:
                return json.dumps({"message": {"content": "Value 42."}}).encode("utf-8")

        with tempfile.TemporaryDirectory() as temporary, patch(
            "book_to_latex.urllib.request.urlopen", return_value=FakeResponse()
        ) as mocked_open:
            image_path = Path(temporary) / "page.png"
            image_path.write_bytes(b"fake-png")
            result = _query_llm(
                page_text="Value 42.",
                page_no=0,
                source_format="PDF",
                page_image_path=image_path,
                style_guide="Use \\Prob for each problem.",
                redraw_graphs=True,
                model="vision-model",
                endpoint=DEFAULT_OLLAMA_ENDPOINT,
                api_key="",
                max_tokens=100,
                temperature=0,
                timeout=10,
                retries=0,
                backoff=1,
                strict_mode=True,
                match_mode="exact",
                log_callback=None,
                preserve_layout=True,
                preserve_color=True,
                page_flow=PAGE_FLOW_COMPACT,
                photo_handling=PHOTO_DESCRIBE,
                page_asset_references=["finished_assets/photo.png"],
            )
            self.assertEqual(result, "Value 42.")
            request = mocked_open.call_args.args[0]
            payload = json.loads(request.data.decode("utf-8"))
            self.assertIs(payload["think"], False)
            self.assertEqual(payload["messages"][1]["images"], ["ZmFrZS1wbmc="])
            self.assertIn("Use \\Prob", payload["messages"][0]["content"])
            self.assertIn("finished_assets/photo.png", payload["messages"][0]["content"])
            self.assertIn("Photo description:", payload["messages"][0]["content"])
            self.assertIn("Semantic visual reconstruction is required", payload["messages"][0]["content"])
            self.assertIn("compact continuous document", payload["messages"][0]["content"])

    def test_ollama_discovery_has_stable_shape(self) -> None:
        from book_to_latex import ollama_connection_info

        info = ollama_connection_info(timeout=1)
        self.assertIn("available", info)
        self.assertIn("models", info)
        if info["available"]:
            self.assertTrue(all("name" in model for model in info["models"]))


class BroadInputTests(unittest.TestCase):
    def test_docx_text_and_table_are_extracted(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            path = Path(temporary) / "book.docx"
            document = Document()
            document.add_heading("Chapter One", level=1)
            document.add_paragraph("The opening paragraph.")
            table = document.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "Name"
            table.cell(0, 1).text = "Value"
            table.cell(1, 0).text = "Alpha"
            table.cell(1, 1).text = "42"
            document.save(path)
            pages, input_format = _read_document_pages(path, 70)
            combined = "\n".join(pages)
            self.assertEqual(input_format, "Word document")
            self.assertIn("Chapter One", combined)
            self.assertIn("Alpha\t42", combined)

    def test_powerpoint_slides_are_kept_as_content_units(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            path = Path(temporary) / "slides.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Project title"
            slide.placeholders[1].text = "First important point"
            presentation.save(path)
            pages, input_format = _read_document_pages(path, 70)
            self.assertEqual(input_format, "PowerPoint presentation")
            self.assertEqual(len(pages), 1)
            self.assertIn("Project title", pages[0])

    def test_excel_sheets_and_values_are_extracted(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            path = Path(temporary) / "data.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Results"
            sheet.append(["Name", "Score"])
            sheet.append(["Ada", 99])
            workbook.save(path)
            pages, input_format = _read_document_pages(path, 70)
            self.assertEqual(input_format, "Excel workbook")
            self.assertIn("[Sheet: Results]", pages[0])
            self.assertIn("Ada\t99", pages[0])

    def test_html_and_unknown_text_extension_are_accepted(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            html_path = root / "article.html"
            html_path.write_text(
                "<html><body><h1>Visible title</h1><script>hidden()</script><p>Article text</p></body></html>",
                encoding="utf-8",
            )
            html_pages, _ = _read_document_pages(html_path, 70)
            self.assertIn("Visible title", html_pages[0])
            self.assertNotIn("hidden", html_pages[0])
            unknown = root / "notes.custom-extension"
            unknown.write_text("Readable content with an unknown extension", encoding="utf-8")
            pages, input_format = _read_document_pages(unknown, 70)
            self.assertEqual(input_format, "text file")
            self.assertIn("Readable content", pages[0])

    def test_unknown_binary_gets_a_clear_error(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            path = Path(temporary) / "random.binary"
            path.write_bytes(b"\x00\x01\x02\x03\xff\x00")
            with self.assertRaisesRegex(ValueError, "unsupported binary"):
                _read_document_pages(path, 70)

    def test_legacy_word_file_through_libreoffice_when_available(self) -> None:
        libreoffice = runtime_capabilities().get("libreoffice")
        if not libreoffice:
            self.skipTest("LibreOffice is not installed")
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            modern_path = root / "legacy-source.docx"
            document = Document()
            document.add_paragraph("Legacy document content")
            document.save(modern_path)
            completed = subprocess.run(
                [
                    str(libreoffice),
                    "--headless",
                    "--convert-to",
                    "doc",
                    "--outdir",
                    str(root),
                    str(modern_path),
                ],
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )
            legacy_path = root / "legacy-source.doc"
            self.assertEqual(completed.returncode, 0, completed.stderr or completed.stdout)
            self.assertTrue(legacy_path.is_file())
            pages, input_format = _read_document_pages(legacy_path, 70)
            self.assertEqual(input_format, "legacy Office document")
            self.assertIn("Legacy document content", "\n".join(pages))


class CliTests(unittest.TestCase):
    def test_cli_failure_returns_nonzero_status(self) -> None:
        completed = subprocess.run(
            [
                sys.executable,
                str(Path(__file__).parents[1] / "book_to_latex.py"),
                "--input",
                "missing-file.txt",
                "--output",
                "unused.tex",
                "--no-llm",
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        self.assertNotEqual(completed.returncode, 0)
        self.assertIn("Conversion failed", completed.stderr)


if __name__ == "__main__":
    unittest.main()
